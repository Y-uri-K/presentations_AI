from __future__ import annotations

import io
import logging
import zipfile
from copy import deepcopy
from dataclasses import dataclass, field
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt

logger = logging.getLogger(__name__)

_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


@dataclass
class TextStyleHint:
    font_size_pt: Optional[float] = None
    rgb: Optional[RGBColor] = None
    bold: Optional[bool] = None


@dataclass
class UserTemplateStyle:
    """Стилистика шаблона без копирования макетов и декора слайдов."""
    theme_files: Dict[str, bytes] = field(default_factory=dict)
    title_text_style: TextStyleHint = field(default_factory=TextStyleHint)
    body_text_style: TextStyleHint = field(default_factory=TextStyleHint)
    accent_rgb: Optional[RGBColor] = None
    dark_rgb: Optional[RGBColor] = None
    light_rgb: Optional[RGBColor] = None
    font_family: Optional[str] = None
    card_fill_rgb: Optional[RGBColor] = None
    background_image: Optional[bytes] = None
    key_elements: List[str] = field(
        default_factory=lambda: ["accent_top_bar", "rounded_cards"]
    )
    title_layout_pattern: str = "title_center"
    title_logo_zone: Optional[str] = None
    title_hero_mode: Optional[str] = None
    content_image_side: str = "right"
    diagram_schemes: List[str] = field(default_factory=list)
    shape_type_counts: Dict[str, int] = field(default_factory=dict)
    palette_hex: List[str] = field(default_factory=list)


def extract_user_template_style(template_bytes: bytes) -> UserTemplateStyle:
    prs = Presentation(io.BytesIO(template_bytes))
    accent_rgb, dark_rgb, light_rgb = _extract_theme_rgb_colors(template_bytes)
    style = UserTemplateStyle(
        theme_files=_extract_theme_files(template_bytes),
        title_text_style=_extract_title_style(prs),
        body_text_style=_extract_body_style(prs),
        accent_rgb=accent_rgb,
        dark_rgb=dark_rgb,
        light_rgb=light_rgb,
        font_family=_extract_font_family_from_theme(template_bytes),
    )
    logger.info(
        "Стиль шаблона (только палитра/шрифты): theme=%s файлов, шрифт=%s",
        len(style.theme_files),
        style.font_family or "—",
    )
    return style


def _extract_theme_files(template_bytes: bytes) -> Dict[str, bytes]:
    files: Dict[str, bytes] = {}
    with zipfile.ZipFile(io.BytesIO(template_bytes)) as archive:
        for name in archive.namelist():
            if name.startswith("ppt/theme/"):
                files[name] = archive.read(name)
    return files


def _parse_rgb(element) -> Optional[RGBColor]:
    if element is None:
        return None
    srgb = element.find("a:srgbClr", _NS)
    if srgb is not None and srgb.get("val"):
        value = srgb.get("val")
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
    return None


def _extract_theme_rgb_colors(
    template_bytes: bytes,
) -> tuple[Optional[RGBColor], Optional[RGBColor], Optional[RGBColor]]:
    accent = None
    dark = None
    light = None
    with zipfile.ZipFile(io.BytesIO(template_bytes)) as archive:
        theme_names = sorted(n for n in archive.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml"))
        if not theme_names:
            return None, None, None
        root = _parse_xml(archive.read(theme_names[0]))
        scheme = root.find(".//a:clrScheme", _NS)
        if scheme is None:
            return None, None, None
        accent = _parse_rgb(scheme.find("a:accent1", _NS))
        dark = _parse_rgb(scheme.find("a:dk1", _NS))
        light = _parse_rgb(scheme.find("a:lt2", _NS)) or _parse_rgb(scheme.find("a:lt1", _NS))
    return accent, dark, light


def _extract_font_family_from_theme(template_bytes: bytes) -> Optional[str]:
    with zipfile.ZipFile(io.BytesIO(template_bytes)) as archive:
        theme_names = sorted(n for n in archive.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml"))
        if not theme_names:
            return None
        root = _parse_xml(archive.read(theme_names[0]))
    for tag in ("majorFont", "minorFont"):
        latin = root.find(f".//a:{tag}/a:latin", _NS)
        if latin is not None and latin.get("typeface"):
            return latin.get("typeface")
    return None


def _parse_xml(data: bytes):
    from lxml import etree

    return etree.fromstring(data)


def _run_rgb(run) -> Optional[RGBColor]:
    try:
        if run.font.color and run.font.color.rgb:
            return run.font.color.rgb
    except AttributeError:
        pass
    return None


def _first_run_style(shape) -> TextStyleHint:
    if not getattr(shape, "has_text_frame", False):
        return TextStyleHint()
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text.strip():
                continue
            size_pt = run.font.size.pt if run.font.size is not None else None
            return TextStyleHint(font_size_pt=size_pt, rgb=_run_rgb(run), bold=run.font.bold)
    return TextStyleHint()


def _extract_title_style(prs: Presentation) -> TextStyleHint:
    for slide in prs.slides:
        for shape in slide.placeholders:
            try:
                if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    hint = _first_run_style(shape)
                    if hint.font_size_pt or hint.rgb:
                        return hint
            except (AttributeError, ValueError):
                continue
        if slide.shapes.title:
            hint = _first_run_style(slide.shapes.title)
            if hint.font_size_pt or hint.rgb:
                return hint
    return TextStyleHint()


def _extract_body_style(prs: Presentation) -> TextStyleHint:
    for slide in prs.slides:
        for shape in slide.placeholders:
            try:
                if shape.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                    hint = _first_run_style(shape)
                    if hint.font_size_pt or hint.rgb:
                        return hint
            except (AttributeError, ValueError):
                continue
    return TextStyleHint()


def merge_theme_files(pptx_bytes: bytes, theme_files: Dict[str, bytes]) -> bytes:
    if not theme_files:
        return pptx_bytes

    input_buffer = io.BytesIO(pptx_bytes)
    output_buffer = io.BytesIO()
    with zipfile.ZipFile(input_buffer, "r") as source, zipfile.ZipFile(
        output_buffer, "w", compression=zipfile.ZIP_DEFLATED
    ) as target:
        for item in source.infolist():
            if item.filename.startswith("ppt/theme/"):
                continue
            target.writestr(item, source.read(item.filename))
        for name, data in theme_files.items():
            target.writestr(name, data)
    return output_buffer.getvalue()

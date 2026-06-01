from __future__ import annotations

import io
import logging
from typing import Dict, Optional

from pptx import Presentation

from app.schemas.slides import PresentationSlides
from app.services.pptx_kinds import SlideBlueprint
from app.services.slide_build_controller import build_slide_controlled
from app.services.slide_templates import SlideVariantPicker
from app.services.template_style_extractor import UserTemplateStyle, merge_theme_files

logger = logging.getLogger(__name__)


def _slide_size_from_template(
    template_bytes: bytes,
    template_file_type: str,
) -> tuple[Optional[int], Optional[int]]:
    if template_file_type != "pptx":
        return None, None
    try:
        source = Presentation(io.BytesIO(template_bytes))
        return source.slide_width, source.slide_height
    except Exception:
        return None, None


def _blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if "blank" in name or "пуст" in name:
            return layout
    return min(prs.slide_layouts, key=lambda layout: len(layout.placeholders))


def _create_presentation(
    template_bytes: bytes,
    template_file_type: str,
) -> Presentation:
    width, height = _slide_size_from_template(template_bytes, template_file_type)
    prs = Presentation()
    if width and height:
        prs.slide_width = width
        prs.slide_height = height
    return prs


def build_pptx_from_template(
    *,
    template_bytes: bytes,
    template_file_type: str,
    slides: PresentationSlides,
    slide_images: Dict[int, bytes],
    user_style: UserTemplateStyle,
) -> bytes:
    """
    Создаёт презентацию с нуля: фон из стиля шаблона → контент → типографика.
  """
    prs = _create_presentation(template_bytes, template_file_type)
    blank_layout = _blank_layout(prs)
    blank_blueprint = SlideBlueprint(
        kind="blank",
        layout=blank_layout,
        layout_name=blank_layout.name or "Blank",
    )

    variant_picker = SlideVariantPicker()

    logger.info(
        "Сборка PPTX: %s слайдов, палитра=%s",
        len(slides.slides),
        ", ".join(user_style.palette_hex[:6]) if user_style.palette_hex else "—",
    )

    for index, slide_spec in enumerate(slides.slides):
        image_bytes = slide_images.get(index)
        slide = prs.slides.add_slide(blank_blueprint.layout)
        has_image = slide_spec.image.source in ("generate", "materials") or index in slide_images
        layout_variant = variant_picker.pick(slide_spec.type, index, has_image=has_image)
        build_slide_controlled(
            slide=slide,
            spec=slide_spec,
            slide_index=index,
            slide_width=prs.slide_width,
            slide_height=prs.slide_height,
            user_style=user_style,
            image_bytes=image_bytes,
            layout_variant=layout_variant,
        )

    output = io.BytesIO()
    prs.save(output)
    built_bytes = output.getvalue()
    return merge_theme_files(built_bytes, user_style.theme_files)

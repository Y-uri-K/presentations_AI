from __future__ import annotations

import io
import logging
from collections import Counter
from dataclasses import dataclass, field
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.template_slide_inspector import _rgb_to_hex, _shape_fill_hex, _slide_background_hex

logger = logging.getLogger(__name__)


@dataclass
class ShapeDetail:
    kind: str
    role: str
    left_pct: float
    top_pct: float
    width_pct: float
    height_pct: float
    fill_hex: Optional[str] = None
    text_preview: str = ""
    auto_shape_name: Optional[str] = None


@dataclass
class SlideGeometryDetail:
    slide_index: int
    layout_pattern: str
    diagram_schemes: List[str] = field(default_factory=list)
    shapes: List[ShapeDetail] = field(default_factory=list)


@dataclass
class TemplateGeometryReport:
    title_slide: SlideGeometryDetail
    slides: List[SlideGeometryDetail]
    shape_type_counts: Dict[str, int] = field(default_factory=dict)
    scheme_counts: Dict[str, int] = field(default_factory=dict)
    recommended_content_image_side: str = "right"
    title_logo_zone: Optional[str] = None
    title_hero_mode: Optional[str] = None


def _pct(value: int, total: int) -> float:
    if not total:
        return 0.0
    return round(value / total, 3)


def _auto_shape_label(shape) -> Optional[str]:
    try:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            return None
        name = str(shape.auto_shape_type).split(".")[-1].lower()
        return name
    except Exception:
        return None


def _guess_role(
    shape,
    *,
    kind: str,
    left_pct: float,
    top_pct: float,
    width_pct: float,
    height_pct: float,
    slide_index: int,
    text: str,
) -> str:
    area = width_pct * height_pct
    if kind == "picture":
        if slide_index == 1 and area < 0.08 and top_pct < 0.2 and left_pct < 0.25:
            return "logo"
        if area > 0.35:
            return "hero_image"
        return "illustration"
    if kind == "line":
        if width_pct > 0.5 and height_pct < 0.05:
            return "connector_horizontal"
        if height_pct > 0.4 and width_pct < 0.05:
            return "connector_vertical"
        return "line"
    if kind == "chart":
        return "chart"
    if kind == "table":
        return "table"
    if height_pct < 0.04 and width_pct > 0.7 and top_pct < 0.08:
        return "accent_bar"
    if text and area > 0.04:
        if slide_index == 1 and top_pct < 0.45 and width_pct > 0.35:
            return "title_text"
        if slide_index == 1 and top_pct > 0.35:
            return "subtitle_text"
        return "body_text"
    if kind in ("auto_shape", "textbox") and area > 0.02:
        if width_pct < 0.35 and height_pct < 0.35:
            return "card"
        return "decorative_shape"
    return "decorative_shape"


def _infer_layout_pattern(slide_index: int, shapes: List[ShapeDetail]) -> str:
    roles = {shape.role for shape in shapes}
    kinds = Counter(shape.kind for shape in shapes)

    if slide_index == 1:
        if "hero_image" in roles and any(s.role == "hero_image" and s.left_pct > 0.45 for s in shapes):
            return "title_hero_right"
        if "hero_image" in roles and any(s.role == "hero_image" and s.left_pct < 0.2 for s in shapes):
            return "title_hero_left"
        if "logo" in roles:
            return "title_logo_top"
        if "title_text" in roles:
            return "title_center"
        return "title_minimal"

    cards = sum(1 for s in shapes if s.role == "card")
    connectors = sum(1 for s in shapes if s.role.startswith("connector"))
    lines = kinds.get("line", 0)

    if kinds.get("table", 0):
        return "table_grid"
    if kinds.get("chart", 0):
        return "chart_diagram"
    if connectors >= 2 or (lines >= 3 and cards >= 2):
        return "timeline_or_flow"
    if cards >= 4:
        return "card_grid"
    if cards == 2 and any(s.left_pct < 0.48 for s in shapes if s.role == "card") and any(
        s.left_pct > 0.5 for s in shapes if s.role == "card"
    ):
        return "two_column"
    if kinds.get("picture", 0) and cards:
        return "image_with_cards"
    return "content_standard"


def _infer_diagram_schemes(layout_pattern: str, shapes: List[ShapeDetail]) -> List[str]:
    schemes: List[str] = []
    if layout_pattern in ("timeline_or_flow",):
        schemes.append("horizontal_timeline")
    if layout_pattern == "card_grid":
        schemes.append("card_grid")
    if layout_pattern == "two_column":
        schemes.append("comparison_columns")
    if layout_pattern == "chart_diagram":
        schemes.append("chart")
    if layout_pattern == "table_grid":
        schemes.append("table")
    if layout_pattern.startswith("title_hero"):
        schemes.append("title_hero")

    auto_names = [s.auto_shape_name for s in shapes if s.auto_shape_name]
    name_blob = " ".join(auto_names).lower()
    if "chevron" in name_blob or "arrow" in name_blob:
        schemes.append("chevron_process")
    if "wave" in name_blob or sum(1 for s in shapes if s.role == "line") >= 4:
        schemes.append("wave_or_connector")
    if sum(1 for s in shapes if s.role == "card") >= 3:
        schemes.append("rounded_cards")
    return sorted(set(schemes))


def _analyze_slide_geometry(slide, slide_index: int, slide_width: int, slide_height: int) -> SlideGeometryDetail:
    shapes: List[ShapeDetail] = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            kind = "picture"
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            kind = "line"
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            kind = "chart"
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            kind = "table"
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            kind = "auto_shape"
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            kind = "textbox"
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            kind = "group"
        elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            kind = "freeform"
        else:
            kind = "shape"

        left_pct = _pct(shape.left, slide_width)
        top_pct = _pct(shape.top, slide_height)
        width_pct = _pct(shape.width, slide_width)
        height_pct = _pct(shape.height, slide_height)
        text = ""
        if getattr(shape, "has_text_frame", False):
            text = (shape.text_frame.text or "").strip()[:80]

        role = _guess_role(
            shape,
            kind=kind,
            left_pct=left_pct,
            top_pct=top_pct,
            width_pct=width_pct,
            height_pct=height_pct,
            slide_index=slide_index,
            text=text,
        )
        shapes.append(
            ShapeDetail(
                kind=kind,
                role=role,
                left_pct=left_pct,
                top_pct=top_pct,
                width_pct=width_pct,
                height_pct=height_pct,
                fill_hex=_shape_fill_hex(shape),
                text_preview=text,
                auto_shape_name=_auto_shape_label(shape),
            )
        )

    layout_pattern = _infer_layout_pattern(slide_index, shapes)
    schemes = _infer_diagram_schemes(layout_pattern, shapes)
    return SlideGeometryDetail(
        slide_index=slide_index,
        layout_pattern=layout_pattern,
        diagram_schemes=schemes,
        shapes=shapes,
    )


def _recommend_image_side(title_slide: SlideGeometryDetail) -> str:
    heroes = [s for s in title_slide.shapes if s.role == "hero_image"]
    if not heroes:
        return "right"
    hero = max(heroes, key=lambda s: s.width_pct * s.height_pct)
    center_x = hero.left_pct + hero.width_pct / 2
    if center_x < 0.42:
        return "right"
    if center_x > 0.58:
        return "left"
    return "right"


def _title_logo_zone(title_slide: SlideGeometryDetail) -> Optional[str]:
    logos = [s for s in title_slide.shapes if s.role == "logo"]
    if not logos:
        return None
    logo = logos[0]
    if logo.left_pct < 0.35 and logo.top_pct < 0.25:
        return "top_left"
    if logo.left_pct > 0.6 and logo.top_pct < 0.25:
        return "top_right"
    return "top_center"


def analyze_template_geometry(
    template_bytes: bytes,
    *,
    max_slides: int = 10,
) -> TemplateGeometryReport:
    prs = Presentation(io.BytesIO(template_bytes))
    width = prs.slide_width or 1
    height = prs.slide_height or 1
    limit = min(len(prs.slides), max_slides)
    slides: List[SlideGeometryDetail] = []

    for index in range(limit):
        slides.append(_analyze_slide_geometry(prs.slides[index], index + 1, width, height))

    title_slide = slides[0] if slides else SlideGeometryDetail(slide_index=1, layout_pattern="title_minimal")

    shape_counter: Counter[str] = Counter()
    scheme_counter: Counter[str] = Counter()
    for slide in slides:
        for shape in slide.shapes:
            shape_counter[shape.kind] += 1
        for scheme in slide.diagram_schemes:
            scheme_counter[scheme] += 1

    hero_mode = None
    if title_slide.layout_pattern.startswith("title_hero"):
        hero_mode = title_slide.layout_pattern.replace("title_", "")

    report = TemplateGeometryReport(
        title_slide=title_slide,
        slides=slides,
        shape_type_counts=dict(shape_counter),
        scheme_counts=dict(scheme_counter),
        recommended_content_image_side=_recommend_image_side(title_slide),
        title_logo_zone=_title_logo_zone(title_slide),
        title_hero_mode=hero_mode,
    )
    logger.info(
        "Геометрия шаблона: титул=%s, схемы=%s, фигуры=%s, картинки контента=%s",
        title_slide.layout_pattern,
        title_slide.diagram_schemes,
        dict(shape_counter),
        report.recommended_content_image_side,
    )
    return report


def geometry_for_polza(report: TemplateGeometryReport) -> str:
    import json

    def shape_row(shape: ShapeDetail) -> dict:
        return {
            "kind": shape.kind,
            "role": shape.role,
            "box": [shape.left_pct, shape.top_pct, shape.width_pct, shape.height_pct],
            "fill": shape.fill_hex,
            "auto": shape.auto_shape_name,
            "text": shape.text_preview[:40],
        }

    payload = {
        "title_slide": {
            "pattern": report.title_slide.layout_pattern,
            "logo_zone": report.title_logo_zone,
            "hero_mode": report.title_hero_mode,
            "schemes": report.title_slide.diagram_schemes,
            "key_shapes": [shape_row(s) for s in report.title_slide.shapes[:20]],
        },
        "shape_type_counts": report.shape_type_counts,
        "scheme_counts": report.scheme_counts,
        "content_image_side_hint": report.recommended_content_image_side,
        "slides_summary": [
            {
                "slide": s.slide_index,
                "pattern": s.layout_pattern,
                "schemes": s.diagram_schemes,
                "shape_roles": sorted({sh.role for sh in s.shapes}),
            }
            for s in report.slides[1:8]
        ],
    }
    return json.dumps(payload, ensure_ascii=False, indent=2)

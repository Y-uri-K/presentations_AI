from __future__ import annotations

import io
from functools import lru_cache
from pathlib import Path

from pptx import Presentation

ASSETS_DIR = Path(__file__).resolve().parent.parent / "assets"
SYSTEM_TEMPLATE_PATH = ASSETS_DIR / "system_template.pptx"


def ensure_system_template() -> Path:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    if SYSTEM_TEMPLATE_PATH.exists():
        return SYSTEM_TEMPLATE_PATH

    prs = Presentation()
    layouts = list(prs.slide_layouts)
    if not layouts:
        prs.save(SYSTEM_TEMPLATE_PATH)
        return SYSTEM_TEMPLATE_PATH

    title_idx = 0
    content_idx = min(1, len(layouts) - 1)
    section_idx = min(2, len(layouts) - 1) if len(layouts) > 2 else content_idx
    picture_idx = min(8, len(layouts) - 1) if len(layouts) > 8 else content_idx

    title_slide = prs.slides.add_slide(layouts[title_idx])
    title_slide.shapes.title.text = "Заголовок презентации"

    section_slide = prs.slides.add_slide(layouts[section_idx])
    section_slide.shapes.title.text = "Раздел"

    content_slide = prs.slides.add_slide(layouts[content_idx])
    content_slide.shapes.title.text = "Заголовок слайда"
    if len(content_slide.placeholders) > 1:
        content_slide.placeholders[1].text = "Пункт списка"

    picture_slide = prs.slides.add_slide(layouts[picture_idx])
    picture_slide.shapes.title.text = "Слайд с изображением"

    prs.save(SYSTEM_TEMPLATE_PATH)
    return SYSTEM_TEMPLATE_PATH


@lru_cache
def get_system_template_bytes() -> bytes:
    path = ensure_system_template()
    return path.read_bytes()


def get_system_presentation() -> Presentation:
    return Presentation(io.BytesIO(get_system_template_bytes()))

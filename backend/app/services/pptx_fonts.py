from __future__ import annotations

from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

DEFAULT_PRESENTATION_FONT = "Times New Roman"


def _sub_element(parent, tag: str, **attrs):
    element = OxmlElement(tag)
    for key, value in attrs.items():
        element.set(key, value)
    parent.append(element)
    return element


def apply_font_to_run(run, font_name: str = DEFAULT_PRESENTATION_FONT) -> None:
    run.font.name = font_name
    r_pr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        element = r_pr.find(qn(f"a:{tag}"))
        if element is None:
            element = _sub_element(r_pr, f"a:{tag}")
        element.set("typeface", font_name)


def apply_font_family_preserve_size(run, font_name: str = DEFAULT_PRESENTATION_FONT) -> None:
    size = run.font.size
    bold = run.font.bold
    italic = run.font.italic
    apply_font_to_run(run, font_name)
    if size is not None:
        run.font.size = size
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic


def apply_font_family_preserve_to_text_frame(
    text_frame,
    font_name: str = DEFAULT_PRESENTATION_FONT,
) -> None:
    for paragraph in text_frame.paragraphs:
        if paragraph.runs:
            for run in paragraph.runs:
                apply_font_family_preserve_size(run, font_name)
        elif paragraph.text:
            run = paragraph.add_run()
            run.text = paragraph.text
            paragraph.text = ""
            apply_font_family_preserve_size(run, font_name)


def apply_font_family_preserve_to_slide(slide, font_name: str = DEFAULT_PRESENTATION_FONT) -> None:
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            apply_font_family_preserve_to_text_frame(shape.text_frame, font_name)
    try:
        notes = slide.notes_slide
        if notes.notes_text_frame:
            apply_font_family_preserve_to_text_frame(notes.notes_text_frame, font_name)
    except Exception:
        pass

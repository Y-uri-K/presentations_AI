from __future__ import annotations

import io
import json
import logging
import re
import zlib
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt
from PIL import Image

from app.ai.types import ChatMessage
from app.schemas.slides import (
    AgendaSlide,
    CardItem,
    CardsSlide,
    ComparisonSlide,
    ConclusionSlide,
    DiagramSlide,
    GoalsSlide,
    KpiSlide,
    PresentationSlides,
    ProblemSlide,
    ProcessStep,
    ProcessSlide,
    ResultItem,
    ResultsSlide,
    SemanticSlide,
    SlideImageSpec,
    TableSlide,
    TableSlideData,
    ThankYouSlide,
    TimelineStep,
    TimelineSlide,
    TitleSlide,
)
from app.services.image_normalizer import normalize_image_for_pptx
from app.services.ai_json import extract_json
from app.services.slide_payload_normalizer import normalize_slide_payload
from app.services.slide_text_limits import clamp_text

logger = logging.getLogger(__name__)
ASSETS_DIR = Path(__file__).resolve().parent.parent / "assets"
DVFU_TEMPLATE_PATH = ASSETS_DIR / "dvfu_template.pptx"
DVFU_TEMPLATE_NAME = "ДВФУ"

_SPACE_RE = re.compile(r"\s+")


@dataclass(frozen=True)
class DVFUSlot:
    slide_number: int
    purpose: str
    preferred_types: tuple[str, ...]
    requires_image: bool = False
    field_contract: str = ""


DVFU_SLOTS: tuple[DVFUSlot, ...] = (
    DVFUSlot(1, "Название проекта и ФИО докладчика", ("title",), field_contract="title/subtitle/speaker, заполняется системой"),
    DVFUSlot(2, "Таймлайн", ("timeline", "process"), field_contract="ровно 7 событий timeline.steps; первые 2 label короткие, все description конкретные"),
    DVFUSlot(3, "Три смысловых списка", ("agenda", "goals", "problem", "cards"), field_contract="ровно 3 блока: короткий heading + 1-2 предложения body"),
    DVFUSlot(4, "Этапы создания", ("process", "timeline"), field_contract="ровно 5 этапов process.steps с содержательными title и description"),
    DVFUSlot(5, "Четыре подтемы", ("cards", "goals", "agenda"), field_contract="ровно 4 карточки/подтемы, без длинных абзацев"),
    DVFUSlot(6, "Таблица", ("table",), field_contract="таблица: 3 смысловые колонки + номер строки, 4-8 строк с реальными критериями"),
    DVFUSlot(7, "Схема из четырёх блоков", ("diagram", "process", "cards"), field_contract="ровно 4 key_points/cards, логическая цепочка"),
    DVFUSlot(8, "Три текстовых блока и изображение", ("cards", "problem", "goals"), True, field_contract="ровно 3 текстовых блока; image может быть materials/generate"),
    DVFUSlot(9, "Чеклист и изображение", ("problem", "goals", "cards"), True, field_contract="ровно 3 пункта чеклиста; image может быть materials/generate"),
    DVFUSlot(10, "Три текстовых блока и изображение", ("cards", "comparison", "results"), True, field_contract="ровно 3 блока: контекст, сравнение/вывод, результат; image может быть materials/generate"),
    DVFUSlot(11, "Два блока и результат", ("results", "comparison", "conclusion"), field_contract="ровно 3 result-блока: 2 промежуточных вывода + 1 итог"),
    DVFUSlot(12, "Схема из четырёх последовательных блоков", ("process", "diagram", "cards"), field_contract="ровно 4 последовательных шага, каждый зависит от предыдущего"),
    DVFUSlot(13, "Схема из трёх последовательных блоков", ("diagram", "process", "cards"), field_contract="ровно 3 этапа: heading + краткое описание"),
    DVFUSlot(14, "Финальный слайд", ("thank_you",), field_contract="фиксированный финальный слайд, заполняется системой"),
)

DVFU_AGENT_BRIEF = """
Каталог ДВФУ-шаблона для выбора типа слайда:
- Слайд 1: ОБЯЗАТЕЛЬНЫЙ титульный, только название проекта, подзаголовок и ФИО докладчика.
- Слайд 2: timeline/process, временная шкала: ровно 7 событий; короткие метки + конкретные описания.
- Слайд 3: agenda/goals/problem/cards, ровно 3 смысловых блока: heading + body.
- Слайд 4: process/timeline, ровно 5 этапов создания или реализации.
- Слайд 5: cards/goals/agenda, ровно 4 подтемы.
- Слайд 6: table, таблица: 3 смысловые колонки + номер строки, 4-8 строк.
- Слайд 7: diagram/process/cards, ровно 4 блока схемы.
- Слайд 8: cards/problem/goals, ровно 3 текстовых блока + изображение.
- Слайд 9: problem/goals/cards, ровно 3 пункта чеклиста + изображение.
- Слайд 10: cards/comparison/results, ровно 3 текстовых блока + изображение.
- Слайд 11: results/comparison/conclusion, ровно 3 результата: два промежуточных и итог.
- Слайд 12: process/diagram/cards, ровно 4 последовательных блока со стрелками.
- Слайд 13: diagram/process/cards, ровно 3 последовательных блока.
- Слайд 14: ОБЯЗАТЕЛЬНЫЙ финальный, не менять.

Итоговая презентация должна быть не более 10 слайдов всего:
титульный ДВФУ + до 8 самых подходящих содержательных макетов + финальный ДВФУ.
Не создавай отдельный слайд под название и не создавай отдельный слайд благодарности:
они уже есть в шаблоне.
Избегай повторов: каждый пункт должен раскрывать новую мысль по теме.
Не используй placeholder-текст: «Введите текст», «ТЕКСТ», «Текст 1», «Слово 1», «Список», «ЗАГОЛОВОК», «—».
"""

_PLACEHOLDER_TEXTS = {
    "введите текст",
    "текст",
    "список",
    "заголовок",
    "заголовок основная мысль",
    "подзаголовок",
    "схема",
    "изображение",
    "номер страницы",
    "год",
    "этап",
    "проект",
    "глобальные рынки",
    "ед.",
    "ед. изм.",
    "текущее",
    "целевое",
    "—",
    "-",
}

_OUTLINE_TOPIC_RE = re.compile(r"^#{1,2}\s+(.+)$", re.MULTILINE)
_GENERIC_RE = re.compile(
    r"^(?:показатель|этап|список|пункт|блок|результат|шаг|текст|слово)\s*\d*$",
    re.IGNORECASE,
)
_PAGE_NUMBER_SHAPE_IDS = {
    2: 8,
    3: 3,
    4: 14,
    5: 6,
    6: 9,
    7: 34,
    8: 13,
    9: 20,
    10: 10,
    11: 15,
    12: 26,
    13: 38,
}


def _compact(value: str | None) -> str:
    return _SPACE_RE.sub(" ", (value or "").replace("\v", " ").strip())


def _is_placeholder(value: str | None) -> bool:
    text = _compact(value)
    return text.casefold() in _PLACEHOLDER_TEXTS or bool(_GENERIC_RE.match(text))


def _is_table_service_value(value: str | None) -> bool:
    text = _compact(value).casefold().strip(".")
    return text in {"ед", "ед изм", "текущее", "целевое", "0", "—", "-"}


def get_dvfu_template_bytes() -> bytes:
    if not DVFU_TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"DVFU template not found: {DVFU_TEMPLATE_PATH}")
    return DVFU_TEMPLATE_PATH.read_bytes()


def _outline_chunks(outline: str) -> list[str]:
    body = re.sub(r"<TITLE>\s*.*?\s*</TITLE>", "", outline or "", flags=re.IGNORECASE | re.DOTALL)
    parts = re.split(r"(?=\n#{1,2}\s)", "\n" + body.strip())
    return [part.strip() for part in parts if part.strip() and part.strip().startswith("#")]


def _chunk_title(chunk: str, fallback: str) -> str:
    match = _OUTLINE_TOPIC_RE.search(chunk)
    return _clean(match.group(1) if match else "", fallback=fallback)


def _chunk_bullets(chunk: str) -> list[str]:
    items: list[str] = []
    for line in chunk.splitlines():
        line = line.strip()
        match = re.match(r"^(?:[-*•]|\d+[.)])\s+(.+)$", line)
        if match:
            items.append(match.group(1).strip())
    return _dedupe(items, fallback_prefix="Ключевая мысль")


def _topic_kind(title: str, bullets: Sequence[str]) -> str:
    text = " ".join([title, *bullets]).casefold()
    has_numeric = bool(re.search(r"\d|%|₽|руб|млн|тыс|квартал|месяц|год", text))
    if any(word in text for word in ("таблиц", "показател", "метрик", "kpi")) and has_numeric:
        return "table"
    if any(word in text for word in ("год", "этап", "срок", "хронолог", "таймлайн", "202", "203")):
        return "timeline"
    if any(word in text for word in ("процесс", "алгоритм", "последователь", "шаг", "создан", "разработ")):
        return "process"
    if any(word in text for word in ("результ", "эффект", "итог", "вывод")):
        return "results"
    if any(word in text for word in ("схем", "архитект", "структур", "связ", "модель")):
        return "diagram"
    return "cards"


def _fallback_kind_for_position(index: int) -> str:
    pattern = ("agenda", "process", "cards", "diagram", "results", "cards", "process", "diagram")
    return pattern[(index - 1) % len(pattern)]


def _fallback_kind_for_position_and_title(index: int, title: str, *, slot_variant: int | None = None) -> str:
    patterns = (
        ("agenda", "process", "cards", "diagram", "results", "cards", "process", "diagram"),
        ("goals", "cards", "process", "diagram", "results", "cards", "timeline", "diagram"),
        ("problem", "goals", "cards", "process", "results", "diagram", "cards", "process"),
    )
    variant = (
        slot_variant
        if slot_variant is not None
        else zlib.crc32(title.encode("utf-8", errors="ignore"))
    ) % len(patterns)
    return patterns[variant][(index - 1) % len(patterns[variant])]


def _ensure_outline_sections(
    sections: list[tuple[str, list[str]]],
    *,
    min_count: int = 8,
) -> list[tuple[str, list[str]]]:
    if len(sections) >= min_count:
        return sections[:min_count]

    titles = [title for title, _ in sections if _clean(title)]
    bullets = [bullet for _, items in sections for bullet in items if _clean(bullet)]
    if not titles:
        titles = ["Содержание презентации"]
    if not bullets:
        bullets = titles

    focus_titles = (
        "Контекст",
        "Цели",
        "Подход",
        "Архитектура",
        "Реализация",
        "Проверка",
        "Риски",
        "Итоги",
    )
    existing = {title.casefold() for title, _ in sections}
    result = list(sections)
    cursor = 0
    while len(result) < min_count:
        focus = focus_titles[len(result) % len(focus_titles)]
        base_title = titles[len(result) % len(titles)]
        title = f"{focus}: {base_title}"
        if title.casefold() in existing:
            title = focus
        group = [bullets[(cursor + offset) % len(bullets)] for offset in range(min(3, len(bullets)))]
        cursor += 2
        existing.add(title.casefold())
        result.append((title, group))
    return result[:min_count]


def slides_from_outline_for_dvfu(outline: str, *, slot_variant: int | None = None) -> PresentationSlides:
    semantic_slides: list[SemanticSlide] = []
    chunks = _outline_chunks(outline)
    expanded: list[tuple[str, list[str]]] = []
    for index, chunk in enumerate(chunks, start=1):
        title = _chunk_title(chunk, f"Раздел {index}")
        bullets = _chunk_bullets(chunk)
        expanded.append((title, bullets))
        if len(bullets) >= 4 and len(expanded) < 8:
            midpoint = max(2, len(bullets) // 2)
            expanded[-1] = (title, bullets[:midpoint])
            expanded.append((f"{title}: продолжение", bullets[midpoint:]))
    expanded = _ensure_outline_sections(expanded, min_count=8)

    for index, (title, bullets) in enumerate(expanded[:8], start=1):
        kind = _topic_kind(title, bullets)
        if kind == "cards":
            kind = _fallback_kind_for_position_and_title(index, title, slot_variant=slot_variant)
        if kind == "table":
            rows = [[bullet, "", "", ""] for bullet in bullets[:8]]
            semantic_slides.append(
                TableSlide(
                    title=title,
                    table=TableSlideData(
                        headers=["Критерий", "Содержание", "Вывод"],
                        rows=rows or [[title, "", ""]],
                    ),
                    image=SlideImageSpec(source="none"),
                )
            )
        elif kind == "timeline":
            steps = [
                TimelineStep(label=_word_label(bullet), description=bullet)
                for step_index, bullet in enumerate(bullets[:5], start=1)
            ]
            while len(steps) < 3:
                steps.append(TimelineStep(label=_word_label(title), description=title))
            semantic_slides.append(TimelineSlide(title=title, steps=steps, image=SlideImageSpec(source="none")))
        elif kind == "process":
            steps = [
                ProcessStep(title=_word_label(bullet), description=bullet)
                for step_index, bullet in enumerate(bullets[:5], start=1)
            ]
            while len(steps) < 3:
                steps.append(ProcessStep(title=_word_label(title), description=title))
            semantic_slides.append(ProcessSlide(title=title, steps=steps, image=SlideImageSpec(source="none")))
        elif kind == "results":
            results = [
                ResultItem(label=_word_label(bullet, fallback="Итог"), value=bullet)
                for result_index, bullet in enumerate(bullets[:4], start=1)
            ]
            while len(results) < 2:
                results.append(ResultItem(label=_word_label(title, fallback="Итог"), value=title))
            semantic_slides.append(ResultsSlide(title=title, results=results, image=SlideImageSpec(source="none")))
        elif kind == "diagram":
            semantic_slides.append(
                DiagramSlide(title=title, key_points=bullets[:4] or [title], image=SlideImageSpec(source="none"))
            )
        elif kind == "agenda":
            items = bullets[:5] or [title]
            while len(items) < 3:
                items.append(title)
            semantic_slides.append(AgendaSlide(title=title, items=items, image=SlideImageSpec(source="none")))
        elif kind == "problem":
            pain_points = bullets[:5] or [title]
            while len(pain_points) < 2:
                pain_points.append(title)
            semantic_slides.append(ProblemSlide(title=title, pain_points=pain_points, image=SlideImageSpec(source="none")))
        elif kind == "goals":
            goals = bullets[:5] or [title]
            while len(goals) < 2:
                goals.append(title)
            semantic_slides.append(GoalsSlide(title=title, goals=goals, image=SlideImageSpec(source="none")))
        else:
            cards = [
                CardItem(title=_word_label(bullet), text=bullet)
                for card_index, bullet in enumerate(bullets[:4], start=1)
            ]
            while len(cards) < 2:
                cards.append(CardItem(title=_word_label(title), text=title))
            semantic_slides.append(CardsSlide(title=title, cards=cards, image=SlideImageSpec(source="none")))

    if not semantic_slides:
        semantic_slides.append(
            CardsSlide(
                title="Содержание презентации",
                cards=[
                    CardItem(title="Тема", text="Ключевые положения презентации"),
                    CardItem(title="Итог", text="Ожидаемые результаты и выводы"),
                ],
                image=SlideImageSpec(source="none"),
            )
        )
    return PresentationSlides(slides=semantic_slides)


def _slide_content_for_refinement(slide: SemanticSlide) -> dict:
    data = slide.model_dump(mode="json")
    data.pop("speaker_notes", None)
    data["image"] = data.get("image") or {"source": "none"}
    return data


def _slot_schema_hint(template_slide_number: int) -> dict:
    if template_slide_number == 2:
        return {
            "type": "timeline",
            "title": "краткая основная мысль",
            "steps": [{"label": "год/этап", "description": "событие"}],
        }
    if template_slide_number in (4, 12):
        return {
            "type": "process",
            "title": "краткая основная мысль",
            "steps": [{"title": "этап", "description": "что происходит"}],
        }
    if template_slide_number == 6:
        return {
            "type": "table",
            "title": "краткая основная мысль",
            "table": {
                "headers": ["Критерий", "Содержание", "Вывод"],
                "rows": [["пример", "содержание", "вывод"]],
            },
        }
    if template_slide_number == 11:
        return {
            "type": "results",
            "title": "краткая основная мысль",
            "results": [
                {"label": "промежуточный вывод", "value": "смысл"},
                {"label": "итог", "value": "смысл"},
            ],
        }
    if template_slide_number == 13:
        return {
            "type": "diagram",
            "title": "краткая основная мысль",
            "key_points": ["ровно 3 последовательных пункта"],
        }
    if template_slide_number in (3, 5, 7, 8, 9, 10):
        return {
            "type": "cards",
            "title": "краткая основная мысль",
            "cards": [{"title": "короткий heading", "text": "1-2 предложения"}],
        }
    return {"type": "cards", "title": "краткая основная мысль"}


def _draft_slide_plan(slides: PresentationSlides, deck_title: str, *, slot_variant: int | None = None) -> dict:
    selected, source_indices, selected_template_indices = _select_slots(
        slides,
        deck_title,
        slot_variant=slot_variant,
    )
    result: list[dict] = []
    for page_number, template_index in enumerate(selected_template_indices, start=1):
        template_slide_number = template_index + 1
        slot = DVFU_SLOTS[template_slide_number - 1]
        source_index = source_indices.get(template_index)
        spec = selected[template_index]
        if template_slide_number == 1:
            content = {"type": "title", "title": deck_title, "subtitle": None}
        elif template_slide_number == 14:
            content = {"type": "thank_you", "title": "Спасибо за внимание"}
        elif spec is not None:
            content = _slide_content_for_refinement(spec)
        else:
            content = {"type": "cards", "title": deck_title, "cards": []}
        result.append(
            {
                "page": page_number,
                "template_slide": template_slide_number,
                "purpose": slot.purpose,
                "allowed_types": list(slot.preferred_types),
                "field_contract": slot.field_contract,
                "expected_json_shape": _slot_schema_hint(template_slide_number),
                "empty_fields_policy": "все пустые строки, cards/text, steps/title/description, results/label/value и key_points обязательны к заполнению агентом",
                "source_slide_index": source_index,
                "content": content,
            }
        )
    return {"slides": result}


def _strip_generic_values(value):
    if isinstance(value, str):
        return "" if _is_placeholder(value) else value.strip()
    if isinstance(value, list):
        cleaned = [_strip_generic_values(item) for item in value]
        return [item for item in cleaned if item not in ("", None, [], {})]
    if isinstance(value, dict):
        return {key: cleaned for key, item in value.items() if (cleaned := _strip_generic_values(item)) not in ("", None, [], {})}
    return value


def _filled_values(title: str, values: Sequence[str], count: int) -> list[str]:
    items = [
        item
        for item in _dedupe(values, fallback_prefix="")
        if item
    ]
    if not items:
        items = [_clean(title, fallback="Ключевая мысль")]
    focus = ("Контекст", "Цель", "Подход", "Реализация", "Проверка", "Риски", "Итог", "Вывод")
    cursor = 0
    while len(items) < count:
        base = items[cursor % len(items)]
        label = focus[len(items) % len(focus)]
        candidate = base if label.casefold() in base.casefold() else f"{label}: {base}"
        if candidate.casefold() in {item.casefold() for item in items}:
            candidate = f"{label}: {_clean(title, fallback=base)}"
        items.append(candidate)
        cursor += 1
    return items[:count]


def _full_text_value(value: str, *, title: str) -> str:
    cleaned = _clean(value)
    if len(cleaned.split()) >= 4:
        return cleaned
    topic = _clean(title, fallback="Тема")
    label = cleaned or _word_label(topic, fallback="Ключевая мысль")
    if label.casefold() in topic.casefold():
        return topic
    return f"{label}: {topic}"


def _generic_cards_from_values(title: str, values: Sequence[str], count: int) -> list[dict]:
    items = _filled_values(title, values, count)
    return [
        {
            "title": _word_label(value, fallback=""),
            "text": _truncate_text_naturally(_full_text_value(value, title=title), 180),
        }
        for idx, value in enumerate(items[:count], start=1)
    ]


def _coerce_content_for_template_slot(content: dict, template_slide_number: int) -> dict:
    content = normalize_slide_payload(content)
    title = _clean(content.get("title"), fallback=f"Слайд {template_slide_number}")
    semantic = PresentationSlides.model_validate({"slides": [content]}).slides[0]
    flat_values = _content_values(semantic, fallback_title=title, max_items=8)

    if template_slide_number == 2:
        values = _filled_values(title, flat_values, 7)
        steps = []
        for idx in range(7):
            value = values[idx]
            steps.append(
                {
                    "label": _word_label(value, fallback=""),
                    "description": _truncate_text_naturally(_full_text_value(value, title=title), 140),
                }
            )
        return {"type": "timeline", "title": title, "steps": steps, "image": content.get("image") or {"source": "none"}}
    if template_slide_number in (4, 12):
        count = 5 if template_slide_number == 4 else 4
        values = _filled_values(title, flat_values, count)
        steps = []
        for idx in range(count):
            value = values[idx]
            steps.append(
                {
                    "title": _word_label(value, fallback=""),
                    "description": _truncate_text_naturally(_full_text_value(value, title=title), 160),
                }
            )
        return {"type": "process", "title": title, "steps": steps, "image": content.get("image") or {"source": "none"}}
    if template_slide_number == 6:
        if content.get("type") == "table" and isinstance(content.get("table"), dict):
            table = content["table"]
            headers = [str(item) for item in table.get("headers", [])[:3]]
            rows = table.get("rows", [])[:8]
        else:
            headers = ["Критерий", "Содержание", "Вывод"]
            rows = [
                [value, "", ""]
                for idx, value in enumerate((flat_values or [title])[:8], start=1)
            ]
        while len(headers) < 3:
            headers.append(["Критерий", "Содержание", "Вывод"][len(headers)])
        normalized_rows = []
        for row in rows[:8]:
            if isinstance(row, list):
                normalized = [str(cell) for cell in row[:3]]
            else:
                normalized = [str(row)]
            while len(normalized) < 3:
                normalized.append("")
            normalized_rows.append(normalized)
        return {
            "type": "table",
            "title": title,
            "table": {"headers": headers[:3], "rows": normalized_rows or [[title, "", ""]]},
            "image": content.get("image") or {"source": "none"},
        }
    if template_slide_number in (3, 8, 9, 10, 13):
        count = 3
        if template_slide_number == 13:
            return {
                "type": "diagram",
                "title": title,
                "key_points": [card["text"] for card in _generic_cards_from_values(title, flat_values or [title], count)],
                "image": content.get("image") or {"source": "none"},
            }
        return {
            "type": "cards",
            "title": title,
            "cards": _generic_cards_from_values(title, flat_values or [title], count),
            "image": content.get("image") or {"source": "none"},
        }
    if template_slide_number in (5, 7):
        return {
            "type": "cards",
            "title": title,
            "cards": _generic_cards_from_values(title, flat_values or [title], 4),
            "image": content.get("image") or {"source": "none"},
        }
    if template_slide_number == 11:
        values = _filled_values(title, flat_values, 3)
        results = []
        for idx in range(3):
            value = values[idx]
            results.append(
                {
                    "label": _word_label(value, fallback=""),
                    "value": _truncate_text_naturally(_full_text_value(value, title=title), 160),
                }
            )
        return {"type": "results", "title": title, "results": results, "image": content.get("image") or {"source": "none"}}
    return content


def _slides_from_refined_payload(payload: dict, fallback: PresentationSlides) -> PresentationSlides:
    raw_items = payload.get("slides") if isinstance(payload, dict) else None
    if not isinstance(raw_items, list):
        return fallback

    normalized: list[dict] = []
    for item in raw_items:
        if not isinstance(item, dict):
            continue
        content = item.get("content") if isinstance(item.get("content"), dict) else item
        content = _strip_generic_values(dict(content))
        if not content.get("type") or content.get("type") in {"title", "thank_you"}:
            continue
        template_slide_number = int(item.get("template_slide") or 0)
        content = normalize_slide_payload(content)
        if template_slide_number:
            content = _coerce_content_for_template_slot(content, template_slide_number)
        normalized.append(content)

    if not normalized:
        return fallback
    try:
        return PresentationSlides.model_validate({"slides": normalized[:8]})
    except Exception as exc:
        logger.warning("ДВФУ refinement JSON не прошёл валидацию: %s", exc)
        return fallback


def _adapt_slide_to_template_slot(slide: SemanticSlide | None, template_slide_number: int) -> SemanticSlide | None:
    if slide is None or template_slide_number in (1, 14):
        return slide
    try:
        content = _coerce_content_for_template_slot(
            slide.model_dump(mode="json"),
            template_slide_number,
        )
        return PresentationSlides.model_validate({"slides": [content]}).slides[0]
    except Exception as exc:
        logger.warning(
            "ДВФУ slot adaptation failed: template_slide=%s, type=%s, error=%s",
            template_slide_number,
            getattr(slide, "type", None),
            exc,
        )
        return slide


def _exception_label(exc: Exception) -> str:
    message = str(exc).strip()
    return f"{type(exc).__name__}: {message}" if message else type(exc).__name__


async def refine_dvfu_slides_with_mimo(
    *,
    slides: PresentationSlides,
    outline: str,
    deck_title: str,
    agent_id: str,
    slot_variant: int | None = None,
    timeout_seconds: int = 75,
) -> PresentationSlides:
    draft = _draft_slide_plan(slides, deck_title, slot_variant=slot_variant)
    prompt = f"""Ты проверяешь JSON заполнения презентации ДВФУ перед рендером.

Тебе дан план презентации и черновой JSON слайдов. Каждый элемент содержит:
- page: итоговый номер страницы;
- template_slide: номер слайда в ДВФУ-шаблоне;
- purpose: назначение макета;
- allowed_types: допустимые semantic-типы;
- field_contract: точное количество и смысл полей, которые реально есть в макете;
- expected_json_shape: пример JSON-структуры, которая лучше всего ложится в макет;
- empty_fields_policy: правило заполнения пустых полей;
- content: текущий контент слайда.

Задача:
1. Верни ТОЛЬКО валидный JSON в том же формате: {{"slides":[{{"page":...,"template_slide":...,"content":...}}]}}.
2. Исправь content так, чтобы каждый блок был заполнен текстом строго по теме и плану.
3. Строго соблюдай field_contract: если макет требует ровно 3 блока — верни 3, если 5 этапов — верни 5, если 7 событий — верни 7.
4. type внутри content должен быть одним из allowed_types и соответствовать expected_json_shape.
5. Убери любые generic/placeholder-фразы: "Показатель 1", "Этап 1", "Список 1", "Пункт 1", "Блок 4", "Введите текст", "ТЕКСТ", "Текст 1", "Слово 1".
6. Не повторяй одинаковый текст в разных блоках и на соседних слайдах.
7. Для table используй реальные строки и заголовки из материалов/плана.
8. Для process/timeline подпиши шаги содержательно, не "Шаг 1".
9. Текст каждого блока должен быть кратким: 1 предложение, до 140 символов.
10. Не оставляй пустых строк в обязательных полях макета. Если в черновике пусто, сформулируй недостающий текст по outline.
11. Для коротких полей "Слово" верни обобщающее слово/категорию, а не начало соседнего текста.
12. Не добавляй титульный и финальный содержательные слайды: они уже фиксированы в шаблоне.

Каталог шаблона:
{DVFU_AGENT_BRIEF}

План презентации:
```md
{outline}
```

Черновой JSON:
```json
{json.dumps(draft, ensure_ascii=False, indent=2)}
```
"""
    logger.info(
        "ДВФУ refinement: агент=%s, timeout=%s c, prompt=%s символов, draft_slides=%s",
        agent_id,
        timeout_seconds,
        len(prompt),
        len(draft["slides"]),
    )
    try:
        import asyncio

        from app.ai.registry import chat_with_agent_resilient

        raw = await asyncio.wait_for(
            chat_with_agent_resilient(agent_id, [ChatMessage(role="user", content=prompt)]),
            timeout=timeout_seconds,
        )
        payload = extract_json(raw)
        return _slides_from_refined_payload(payload, slides)
    except Exception as exc:
        logger.warning("ДВФУ refinement через %s пропущен: %s", agent_id, _exception_label(exc))
        return slides


def _clean(value: str | None, *, fallback: str = "") -> str:
    text = _compact(value)
    if _is_placeholder(text):
        text = ""
    return text or fallback


def _dedupe(values: Iterable[str], *, fallback_prefix: str = "Пункт") -> list[str]:
    result: list[str] = []
    seen: set[str] = set()
    for value in values:
        cleaned = _clean(value)
        if not cleaned:
            continue
        key = cleaned.casefold()
        if key in seen:
            continue
        seen.add(key)
        result.append(cleaned)
    if not result:
        result.append(fallback_prefix)
    return result


def _text_capacity(shape) -> int:
    width = max(1, int(shape.width))
    height = max(1, int(shape.height))
    return max(70, min(520, int((width * height) / 28_000_000)))


def _joined(parts: Iterable[str], *, separator: str = "\n") -> str:
    return separator.join(item for item in (_clean(part) for part in parts) if item)


def _short_heading(text: str, *, fallback: str) -> str:
    words = _clean(text, fallback=fallback).split()
    return clamp_text(" ".join(words[:4]) or fallback, 42, ellipsis=False)


def _word_label(text: str, *, fallback: str = "") -> str:
    cleaned = _clean(text)
    lowered = cleaned.casefold()
    rules = (
        (("цель", "задач"), "Цель"),
        (("огранич", "барьер", "риск", "проблем"), "Риски"),
        (("пользовател", "сценар"), "Сценарии"),
        (("критер", "метрик", "показател", "kpi"), "Критерии"),
        (("анализ", "исслед"), "Анализ"),
        (("требован",), "Требования"),
        (("проектир", "структур"), "Структура"),
        (("разработ", "интерфейс"), "Разработка"),
        (("тест", "провер"), "Проверка"),
        (("запуск", "сопровожд"), "Запуск"),
        (("данн", "sql", "баз"), "Данные"),
        (("архитект",), "Архитектура"),
        (("эффект", "результ", "итог", "вывод"), "Итог"),
        (("качество",), "Качество"),
        (("скорост", "время"), "Скорость"),
    )
    for needles, label in rules:
        if any(needle in lowered for needle in needles):
            return label
    words = re.findall(r"[A-Za-zА-Яа-яЁё0-9]+", cleaned)
    if not words:
        return fallback
    return clamp_text(" ".join(words[:2]), 28, ellipsis=False)


def _first_sentence(text: str) -> str:
    cleaned = _clean(text)
    match = re.search(r"^(.+?[.!?])(?:\s|$)", cleaned)
    if match:
        return match.group(1).strip()
    parts = re.split(r"[;:]\s+", cleaned)
    return parts[0].strip() if parts else cleaned


def _truncate_text_naturally(text: str, max_chars: int) -> str:
    cleaned = _clean(text)
    if len(cleaned) <= max_chars:
        return cleaned

    sentences = re.findall(r".+?(?:[.!?](?=\s|$)|$)", cleaned)
    result = ""
    for sentence in sentences:
        candidate = _compact(f"{result} {sentence}" if result else sentence)
        if len(candidate) <= max_chars:
            result = candidate
        else:
            break
    if result:
        return result

    window = cleaned[: max(1, max_chars)].rstrip()
    for separator in (". ", "! ", "? ", "; ", ": ", ", "):
        index = window.rfind(separator)
        if index >= max_chars * 0.45:
            return window[: index + 1].strip()

    space_index = window.rfind(" ")
    if space_index >= max_chars * 0.45:
        return window[:space_index].rstrip(" ,;:") + "."
    return window.rstrip(" ,;:") + "."


def _find_shape(slide, shape_id: int):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    return None


def _set_text(shape, text: str, *, max_chars: int = 260) -> None:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return

    capacity = min(max_chars, _text_capacity(shape))
    raw = _clean(text)
    value = raw
    if len(raw) > capacity:
        max_bottom = 6_130_000
        available = max_bottom - int(shape.top)
        if available > int(shape.height):
            shape.height = min(int(shape.height * 1.45), available)
            capacity = min(max_chars, _text_capacity(shape))
        if len(raw) > capacity:
            one_sentence = _first_sentence(raw)
            value = _truncate_text_naturally(one_sentence, capacity)
        else:
            value = raw
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    lines = value.split("\n") if value else [""]
    first = text_frame.paragraphs[0]
    first.text = lines[0]
    for line in lines[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = line

    font_size = None
    total = len(value)
    if total > 280:
        font_size = Pt(7)
    elif total > 200:
        font_size = Pt(8)
    elif total > 130:
        font_size = Pt(9)
    elif total > 80:
        font_size = Pt(11)

    if font_size is not None:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = font_size


def _clear_text(shape) -> None:
    _set_text(shape, "")


def _set_by_id(slide, shape_id: int, text: str, *, max_chars: int = 260) -> None:
    _set_text(_find_shape(slide, shape_id), text, max_chars=max_chars)


def _set_text_preserve_style(shape, text: str, *, max_chars: int = 260) -> None:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    value = _truncate_text_naturally(text, max_chars)
    paragraphs = list(shape.text_frame.paragraphs)
    if not paragraphs:
        return
    first = paragraphs[0]
    if first.runs:
        first.runs[0].text = value
        for run in first.runs[1:]:
            run.text = ""
    else:
        first.text = value
    for paragraph in paragraphs[1:]:
        for run in paragraph.runs:
            run.text = ""


def _set_by_id_preserve_style(slide, shape_id: int, text: str, *, max_chars: int = 260) -> None:
    _set_text_preserve_style(_find_shape(slide, shape_id), text, max_chars=max_chars)


def _set_page_number(slide, template_slide_number: int, page_number: int) -> None:
    shape_id = _PAGE_NUMBER_SHAPE_IDS.get(template_slide_number)
    if shape_id is not None:
        _set_by_id(slide, shape_id, str(page_number), max_chars=8)
        return
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if _compact(shape.text).casefold() in {"номер страницы"} or _compact(shape.text).isdigit():
            _set_text(shape, str(page_number), max_chars=8)
            return


def _cleanup_placeholders(slide) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if _is_placeholder(shape.text):
            _clear_text(shape)


def _insert_image(slide, shape_id: int, image_bytes: bytes | None) -> None:
    placeholder = _find_shape(slide, shape_id)
    if placeholder is None:
        return

    if not image_bytes:
        _clear_text(placeholder)
        return

    normalized = normalize_image_for_pptx(image_bytes)
    if not normalized:
        _clear_text(placeholder)
        return

    box_w = int(placeholder.width)
    box_h = int(placeholder.height)
    pic_w, pic_h = _fit_picture_size(normalized, box_w, box_h)
    left = int(placeholder.left) + max(0, (box_w - int(pic_w)) // 2)
    top = int(placeholder.top) + max(0, (box_h - int(pic_h)) // 2)
    _clear_text(placeholder)
    slide.shapes.add_picture(io.BytesIO(normalized), left, top, width=int(pic_w), height=int(pic_h))


def _fit_picture_size(image_bytes: bytes, max_width: int, max_height: int) -> tuple[int, int]:
    with Image.open(io.BytesIO(image_bytes)) as image:
        width, height = image.size
    if width <= 0 or height <= 0:
        return max_width, max_height
    scale = min(max_width / width, max_height / height)
    return max(1, int(width * scale)), max(1, int(height * scale))


def _flatten_slide(slide: SemanticSlide) -> list[str]:
    items: list[str] = []
    if isinstance(slide, TitleSlide):
        items.extend([slide.title, slide.subtitle or ""])
    elif isinstance(slide, AgendaSlide):
        items.extend([slide.title, *slide.items])
    elif isinstance(slide, ProblemSlide):
        items.extend([slide.title, *slide.pain_points])
    elif isinstance(slide, GoalsSlide):
        items.extend([slide.title, *slide.goals])
    elif isinstance(slide, KpiSlide):
        items.append(slide.title)
        for metric in slide.metrics:
            items.append(_joined([metric.value, metric.label, metric.note or ""], separator=" — "))
    elif isinstance(slide, CardsSlide):
        items.append(slide.title)
        for card in slide.cards:
            items.append(_card_text(card))
    elif isinstance(slide, ComparisonSlide):
        items.extend([slide.title, slide.left.heading, *slide.left.points, slide.right.heading, *slide.right.points])
    elif isinstance(slide, TimelineSlide):
        items.append(slide.title)
        for step in slide.steps:
            items.append(_joined([step.label, step.description or ""], separator=" — "))
    elif isinstance(slide, ProcessSlide):
        items.append(slide.title)
        for step in slide.steps:
            items.append(_joined([step.title, step.description or ""], separator=" — "))
    elif isinstance(slide, ResultsSlide):
        items.extend([slide.title, slide.summary or ""])
        for result in slide.results:
            items.append(_result_text(result))
    elif isinstance(slide, TableSlide):
        items.extend([slide.title, *slide.table.headers])
        for row in slide.table.rows:
            items.append(" — ".join(row))
    elif isinstance(slide, DiagramSlide):
        items.extend([slide.title, slide.caption or "", *slide.key_points])
    elif isinstance(slide, ConclusionSlide):
        items.extend([slide.title, *slide.takeaways])
    elif isinstance(slide, ThankYouSlide):
        items.extend([slide.title, slide.subtitle or "", slide.contact or ""])
    return [item for item in (_clean(item) for item in items) if item]


def _content_values(slide: SemanticSlide | None, *, fallback_title: str, max_items: int = 8) -> list[str]:
    if slide is None:
        return []

    values: list[str] = []
    if isinstance(slide, AgendaSlide):
        values = list(slide.items)
    elif isinstance(slide, ProblemSlide):
        values = list(slide.pain_points)
    elif isinstance(slide, GoalsSlide):
        values = list(slide.goals)
    elif isinstance(slide, CardsSlide):
        for card in slide.cards:
            text = _clean(card.text)
            title = _clean(card.title)
            if text and title and title.casefold() not in text.casefold():
                values.append(f"{title}: {text}")
            else:
                values.append(text or title)
    elif isinstance(slide, ComparisonSlide):
        values = [
            _joined([slide.left.heading, *slide.left.points], separator=": "),
            _joined([slide.right.heading, *slide.right.points], separator=": "),
        ]
    elif isinstance(slide, TimelineSlide):
        values = [
            _joined([step.label, step.description or ""], separator=": ")
            for step in slide.steps
        ]
    elif isinstance(slide, ProcessSlide):
        values = [
            _joined([step.title, step.description or ""], separator=": ")
            for step in slide.steps
        ]
    elif isinstance(slide, ResultsSlide):
        values = [_result_text(result) for result in slide.results]
        if slide.summary:
            values.append(slide.summary)
    elif isinstance(slide, TableSlide):
        values = [" — ".join(cell for cell in row if _clean(cell)) for row in slide.table.rows]
    elif isinstance(slide, DiagramSlide):
        values = list(slide.key_points)
        if slide.caption:
            values.append(slide.caption)
    elif isinstance(slide, ConclusionSlide):
        values = list(slide.takeaways)
    else:
        values = _flatten_slide(slide)

    cleaned = _dedupe(values, fallback_prefix=fallback_title)
    if cleaned == [fallback_title]:
        cleaned = []
    title = _clean(getattr(slide, "title", None))
    if title:
        cleaned = [value for value in cleaned if value.casefold() != title.casefold()]
    return cleaned[:max_items]


def _card_text(card: CardItem) -> str:
    return _joined([card.title, card.text, card.highlight or ""], separator="\n")


def _result_text(result: ResultItem) -> str:
    value = _clean(result.value)
    label = _clean(result.label)
    if value and label and label.casefold() in value.casefold():
        return value
    return _joined([value, label], separator="\n")


def _title(slide: SemanticSlide | None, fallback: str) -> str:
    return _clean(getattr(slide, "title", None), fallback=fallback)


def _items(slide: SemanticSlide | None, count: int, *, fallback_title: str) -> list[str]:
    title = _title(slide, fallback_title)
    values = _filled_values(title, _content_values(slide, fallback_title=fallback_title, max_items=count), count)
    return [
        _truncate_text_naturally(_full_text_value(value, title=title), 220)
        for idx, value in enumerate(values[:count])
    ]


def _steps(slide: SemanticSlide | None, count: int) -> list[tuple[str, str]]:
    if isinstance(slide, TimelineSlide):
        raw = [(step.label, step.description or "") for step in slide.steps]
    elif isinstance(slide, ProcessSlide):
        raw = [(step.title, step.description or "") for step in slide.steps]
    elif isinstance(slide, CardsSlide):
        raw = [
            (
                _word_label(card.text or card.title, fallback=""),
                _full_text_value(card.text or card.title, title=slide.title),
            )
            for card in slide.cards
        ]
    else:
        values = _items(slide, count, fallback_title="Этап")
        raw = [(_word_label(value, fallback=""), value) for value in values]

    while len(raw) < count:
        raw.append(("", ""))
    return [
        (
            clamp_text(_clean(label), 40),
            _truncate_text_naturally(_full_text_value(description, title=_title(slide, "Тема")), 180),
        )
        for idx, (label, description) in enumerate(raw[:count], start=1)
    ]


def _slot_score(slide_type: str, slot: DVFUSlot) -> int:
    if slide_type in slot.preferred_types:
        return len(slot.preferred_types) - slot.preferred_types.index(slide_type) + 10
    return 0


def _best_slot_for_slide(
    *,
    source: SemanticSlide,
    candidates: Sequence[DVFUSlot],
    used_slide_types: set[str],
    has_image: bool,
    last_slot_number: int,
) -> DVFUSlot | None:
    if not candidates:
        return None

    image_candidates = [slot for slot in candidates if has_image and slot.requires_image]
    if image_candidates:
        candidates = image_candidates

    # ДВФУ-шаблон линейный: если рано занять дальний макет (например, схему 13),
    # последующие содержательные слайды уже некуда будет поставить.
    nearby_candidates = [slot for slot in candidates if slot.slide_number <= last_slot_number + 3]
    if nearby_candidates and not image_candidates:
        candidates = nearby_candidates

    def score(slot: DVFUSlot) -> tuple[int, int, int]:
        match_score = _slot_score(source.type, slot)
        diversity_bonus = 4 if source.type not in used_slide_types else 0
        image_bonus = 30 if has_image and slot.requires_image else 0
        distance_penalty = max(0, slot.slide_number - last_slot_number - 1) * 8
        return (
            match_score + diversity_bonus + image_bonus - distance_penalty,
            -slot.slide_number,
            -len(slot.preferred_types),
        )

    return max(candidates, key=score)


def _content_slot_numbers(
    deck_title: str,
    *,
    has_image_slots: bool,
    slot_variant: int | None = None,
) -> tuple[int, ...]:
    if has_image_slots:
        variants = (
            (2, 3, 4, 5, 8, 9, 10, 11),
            (2, 4, 5, 7, 8, 9, 10, 13),
            (3, 4, 6, 7, 8, 9, 10, 12),
        )
    else:
        variants = (
            (2, 3, 4, 5, 6, 7, 11, 12),
            (2, 4, 5, 6, 7, 11, 12, 13),
            (3, 4, 5, 6, 7, 11, 12, 13),
        )
    variant = (
        slot_variant
        if slot_variant is not None
        else zlib.crc32(deck_title.encode("utf-8", errors="ignore"))
    ) % len(variants)
    return variants[variant]


def _slide_wants_image(slide: SemanticSlide) -> bool:
    return bool(getattr(slide, "image", None) and slide.image.source != "none")


def _source_has_available_image(
    source_index: int,
    slide: SemanticSlide,
    slide_images: dict[int, bytes] | None,
) -> bool:
    if slide_images is not None:
        return source_index in slide_images
    return _slide_wants_image(slide)


def _select_slots(
    slides: PresentationSlides,
    deck_title: str,
    *,
    slide_images: dict[int, bytes] | None = None,
    max_total_slides: int = 10,
    slot_variant: int | None = None,
) -> tuple[list[SemanticSlide | None], dict[int, int], list[int]]:
    selected: list[SemanticSlide | None] = [None] * 14
    source_indices: dict[int, int] = {}
    selected_template_indices = [0]
    used_slots: set[int] = set()
    used_slide_types: set[str] = set()
    content_limit = max(0, max_total_slides - 2)

    selected[0] = TitleSlide(title=deck_title, subtitle=None)
    selected[13] = ThankYouSlide(title="Спасибо за внимание")

    content_slides = [
        (idx, slide)
        for idx, slide in enumerate(slides.slides)
        if slide.type not in {"title", "thank_you", "conclusion"}
    ][:content_limit]
    while len(content_slides) < content_limit:
        synthetic_index = len(content_slides)
        content_slides.append(
            (
                synthetic_index,
                CardsSlide(
                    title=f"Раздел {synthetic_index + 1}: {deck_title}",
                    cards=[],
                    image=SlideImageSpec(source="none"),
                ),
            )
        )

    image_source_count = sum(
        1
        for idx, slide in content_slides
        if _source_has_available_image(idx, slide, slide_images)
    )
    has_image_slots = image_source_count > 0
    route = list(_content_slot_numbers(deck_title, has_image_slots=has_image_slots, slot_variant=slot_variant))
    if has_image_slots:
        image_slots_to_keep = {
            slide_number
            for slide_number in route
            if DVFU_SLOTS[slide_number - 1].requires_image
        }
        image_slots_to_keep = set(sorted(image_slots_to_keep)[:image_source_count])
        route = [
            slide_number
            for slide_number in route
            if not DVFU_SLOTS[slide_number - 1].requires_image or slide_number in image_slots_to_keep
        ]
        for slot in DVFU_SLOTS:
            if len(route) >= content_limit:
                break
            if not (2 <= slot.slide_number <= 13) or slot.requires_image or slot.slide_number in route:
                continue
            route.append(slot.slide_number)
        route = sorted(route)[:content_limit]
    allowed_slot_numbers = set(route)
    content_slots = [
        slot
        for slot in DVFU_SLOTS
        if 2 <= slot.slide_number <= 13 and slot.slide_number in allowed_slot_numbers
    ]

    remaining = list(content_slides)
    for slot_index, slot in enumerate(content_slots):
        image_slots_left = sum(1 for candidate in content_slots[slot_index:] if candidate.requires_image)
        image_sources_left = [
            item for item in remaining if _source_has_available_image(item[0], item[1], slide_images)
        ]
        if slot.requires_image and not image_sources_left:
            continue
        if slot.requires_image and image_sources_left:
            candidates = image_sources_left
        elif not slot.requires_image and len(image_sources_left) >= image_slots_left:
            candidates = [item for item in remaining if item not in image_sources_left] or remaining
        else:
            candidates = remaining
        if not candidates:
            break

        def source_score(item: tuple[int, SemanticSlide]) -> tuple[int, int, int]:
            source_index, source = item
            has_image = _source_has_available_image(source_index, source, slide_images)
            image_score = 30 if has_image and slot.requires_image else -30 if has_image and not slot.requires_image else 0
            type_score = _slot_score(source.type, slot)
            diversity_bonus = 4 if source.type not in used_slide_types else 0
            return (image_score + type_score + diversity_bonus, -source_index, -len(source.type))

        source_index, source = max(candidates, key=source_score)
        remaining.remove((source_index, source))
        template_index = slot.slide_number - 1
        selected[template_index] = source
        source_indices[template_index] = source_index
        used_slots.add(slot.slide_number)
        used_slide_types.add(source.type)
        selected_template_indices.append(template_index)

    selected_template_indices.append(13)
    return selected, source_indices, selected_template_indices[:max_total_slides]


def _fill_table(slide, spec: SemanticSlide | None) -> None:
    table_shape = next((shape for shape in slide.shapes if getattr(shape, "has_table", False)), None)
    if table_shape is None:
        return

    table = table_shape.table
    data_columns = max(1, len(table.columns) - 1)
    default_headers = ["Критерий", "Содержание", "Вывод"][:data_columns]
    if isinstance(spec, TableSlide):
        headers = spec.table.headers[:data_columns]
        rows = spec.table.rows[:8]
    else:
        headers = default_headers
        values = _items(spec, 8, fallback_title="Показатель")
        rows = [[value, *[""] * (data_columns - 1)] for value in values]

    while len(headers) < data_columns:
        headers.append(default_headers[len(headers)] if len(headers) < len(default_headers) else f"Колонка {len(headers) + 1}")
    headers = [
        default_headers[idx] if idx < len(default_headers) else f"Колонка {idx + 1}"
        if _is_table_service_value(header)
        else header
        for idx, header in enumerate(headers[:data_columns])
    ]

    table.cell(0, 0).text = "№"
    for col, header in enumerate(headers[:data_columns], start=1):
        table.cell(0, col).text = clamp_text(_clean(header, fallback=f"Колонка {col}"), 32)

    for row_idx in range(1, len(table.rows)):
        if row_idx - 1 >= len(rows):
            table.cell(row_idx, 0).text = ""
            for col_idx in range(1, len(table.columns)):
                table.cell(row_idx, col_idx).text = ""
            continue
        source = list(rows[row_idx - 1]) if row_idx - 1 < len(rows) else []
        while len(source) < len(headers):
            source.append("")
        table.cell(row_idx, 0).text = str(row_idx)
        for col_idx in range(1, len(table.columns)):
            value = source[col_idx - 1] if col_idx - 1 < len(source) else ""
            if _is_table_service_value(value):
                value = ""
            table.cell(row_idx, col_idx).text = _truncate_text_naturally(value, 80)


def _result_items(spec: SemanticSlide | None) -> list[str]:
    if isinstance(spec, ResultsSlide):
        values = [_result_text(result) for result in spec.results]
        if spec.summary:
            values.append(spec.summary)
        return _dedupe(values, fallback_prefix="Результат")
    if isinstance(spec, ConclusionSlide):
        return _dedupe(spec.takeaways, fallback_prefix="Результат")
    if isinstance(spec, ComparisonSlide):
        return _dedupe(
            [
                _joined([spec.left.heading, *spec.left.points], separator="\n"),
                _joined([spec.right.heading, *spec.right.points], separator="\n"),
                "Итог: " + _joined([*(spec.left.points[:1]), *(spec.right.points[:1])], separator="; "),
            ],
            fallback_prefix="Результат",
        )
    return _items(spec, 3, fallback_title="Результат")


def _fill_slide(
    prs,
    slide_index: int,
    spec: SemanticSlide | None,
    *,
    image_bytes: bytes | None,
    page_number: int,
) -> None:
    slide = prs.slides[slide_index]
    slide_number = slide_index + 1
    if slide_number != 14:
        _set_page_number(slide, slide_number, page_number)

    title = _title(spec, f"Слайд {slide_number}")

    if slide_number == 1:
        return
    if slide_number == 2:
        _set_by_id(slide, 29, title)
        for shape_id, (label, description) in zip((66, 67), _steps(spec, 2)):
            _set_by_id(slide, shape_id, label, max_chars=24)
        for shape_id, (_, description) in zip((37, 2, 3, 4, 5, 6, 7), _steps(spec, 7)):
            _set_by_id(slide, shape_id, description, max_chars=120)
    elif slide_number == 3:
        _set_by_id(slide, 29, title)
        bodies = _items(spec, 3, fallback_title="Список")
        for shape_id, body in zip((63, 68, 70), bodies):
            _set_by_id(slide, shape_id, _word_label(body), max_chars=32)
        for shape_id, body in zip((61, 67, 69), bodies):
            _set_by_id(slide, shape_id, body, max_chars=220)
    elif slide_number == 4:
        _set_by_id(slide, 2, title)
        steps = _steps(spec, 5)
        for shape_id, (label, _) in zip((96, 97, 98, 99, 89), steps):
            _set_by_id(slide, shape_id, label, max_chars=32)
        for shape_id, (_, description) in zip((3, 4, 6, 5, 7), steps):
            _set_by_id(slide, shape_id, description, max_chars=220)
        for shape_id in (88, 90, 114):
            _clear_text(_find_shape(slide, shape_id))
    elif slide_number == 5:
        _set_by_id(slide, 93, title)
        for shape_id, body in zip((63, 65, 66, 71), _items(spec, 4, fallback_title="Подтема")):
            _set_by_id(slide, shape_id, body, max_chars=180)
    elif slide_number == 6:
        _set_by_id(slide, 93, title)
        _fill_table(slide, spec)
    elif slide_number == 7:
        _set_by_id(slide, 3, title)
        _set_by_id(slide, 33, _word_label(title, fallback="Схема"), max_chars=42)
        for shape_id, body in zip((5, 14, 16, 10), _items(spec, 4, fallback_title="Блок")):
            _set_by_id(slide, shape_id, body, max_chars=180)
    elif slide_number == 8:
        _set_by_id(slide, 6, title)
        for shape_id, body in zip((3, 4, 5), _items(spec, 3, fallback_title="Блок")):
            _set_by_id(slide, shape_id, body, max_chars=150)
        _insert_image(slide, 10, image_bytes)
    elif slide_number == 9:
        _set_by_id(slide, 15, title)
        for shape_id, body in zip((18, 5, 19), _items(spec, 3, fallback_title="Пункт")):
            _set_by_id(slide, shape_id, body, max_chars=160)
        _insert_image(slide, 7, image_bytes)
    elif slide_number == 10:
        _set_by_id(slide, 15, title)
        for shape_id, body in zip((4, 5, 8), _items(spec, 3, fallback_title="Блок")):
            _set_by_id(slide, shape_id, body, max_chars=150)
        _insert_image(slide, 7, image_bytes)
    elif slide_number == 11:
        _set_by_id(slide, 6, title)
        items = _result_items(spec)
        while len(items) < 3:
            items.append("")
        _set_by_id(slide, 7, items[0], max_chars=240)
        _set_by_id(slide, 8, items[1], max_chars=240)
        _set_by_id(slide, 9, items[2], max_chars=240)
    elif slide_number == 12:
        _set_by_id(slide, 3, title)
        _set_by_id(slide, 16, _word_label(title, fallback="Схема"), max_chars=42)
        for shape_id, body in zip((17, 18, 19, 20), _items(spec, 4, fallback_title="Этап")):
            _set_by_id(slide, shape_id, body, max_chars=160)
    elif slide_number == 13:
        _set_by_id(slide, 13, title)
        _set_by_id(slide, 35, _word_label(title, fallback="Схема"), max_chars=42)
        items = _items(spec, 3, fallback_title="Блок")
        for shape_id, body in zip((3, 5, 28), items):
            _set_by_id(slide, shape_id, _word_label(body), max_chars=32)
        for shape_id, body in zip((25, 26, 34), items):
            _set_by_id(slide, shape_id, body, max_chars=120)

    _cleanup_placeholders(slide)


def _delete_unselected_slides(prs, keep_indices: Sequence[int]) -> None:
    keep = set(keep_indices)
    slide_id_list = prs.slides._sldIdLst
    for index in reversed(range(len(prs.slides))):
        if index in keep:
            continue
        slide_id = slide_id_list[index]
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def _renumber_final_slides(prs, template_indices: Sequence[int]) -> None:
    for page_number, template_index in enumerate(template_indices, start=1):
        template_slide_number = template_index + 1
        if template_slide_number == 14:
            continue
        _set_page_number(prs.slides[page_number - 1], template_slide_number, page_number)


def _build_selected_slides(
    selected: Sequence[SemanticSlide | None],
    selected_template_indices: Sequence[int],
    deck_title: str,
) -> PresentationSlides:
    result: list[SemanticSlide] = []
    fallback = TitleSlide(title=deck_title, subtitle=None)
    for template_index in selected_template_indices:
        slide = selected[template_index]
        if slide is not None:
            result.append(slide)
        elif template_index == 13:
            result.append(ThankYouSlide(title="Спасибо за внимание"))
        else:
            result.append(fallback)
    return PresentationSlides(slides=result)


def build_dvfu_pptx(
    *,
    template_bytes: bytes,
    slides: PresentationSlides,
    slide_images: dict[int, bytes],
    deck_title: str,
    presenter_name: str | None,
    subtitle: str | None = None,
    slot_variant: int | None = None,
) -> tuple[bytes, PresentationSlides]:
    prs = Presentation(io.BytesIO(template_bytes))
    if len(prs.slides) < 14:
        raise ValueError("ДВФУ-шаблон должен содержать 14 слайдов")

    selected, source_indices, selected_template_indices = _select_slots(
        slides,
        deck_title,
        slide_images=slide_images,
        slot_variant=slot_variant,
    )
    selected = [
        _adapt_slide_to_template_slot(slide, index + 1)
        for index, slide in enumerate(selected)
    ]
    slide1 = prs.slides[0]
    _set_by_id_preserve_style(slide1, 10, deck_title, max_chars=120)
    _set_by_id_preserve_style(slide1, 19, subtitle or "Проектная презентация", max_chars=80)
    _set_by_id_preserve_style(slide1, 9, presenter_name or "Докладчик", max_chars=80)

    for page_number, template_index in enumerate(selected_template_indices, start=1):
        if template_index in (0, 13):
            continue
        source_index = source_indices.get(template_index)
        image_bytes = slide_images.get(source_index) if source_index is not None else None
        _fill_slide(
            prs,
            template_index,
            selected[template_index],
            image_bytes=image_bytes,
            page_number=page_number,
        )

    _delete_unselected_slides(prs, selected_template_indices)
    _renumber_final_slides(prs, selected_template_indices)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue(), _build_selected_slides(selected, selected_template_indices, deck_title)

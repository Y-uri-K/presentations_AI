from __future__ import annotations

import io
import logging
import zipfile
from collections import Counter
from typing import List, Optional, Tuple

import fitz
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.pdf_template_style_extractor import _color_to_int, _palette_from_page
from app.services.template_slide_inspector import _rgb_to_hex, _shape_fill_hex, _slide_background_hex
from app.services.template_style_extractor import TextStyleHint, UserTemplateStyle, _parse_rgb

logger = logging.getLogger(__name__)

_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}


def _parse_xml(data: bytes):
    from lxml import etree

    return etree.fromstring(data)


def _luminance(rgb: RGBColor) -> float:
    try:
        value = _color_to_int(rgb)
    except (TypeError, ValueError):
        return 0.5
    r, g, b = (value >> 16) & 0xFF, (value >> 8) & 0xFF, value & 0xFF
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255.0


def _saturation(rgb: RGBColor) -> float:
    try:
        value = _color_to_int(rgb)
    except (TypeError, ValueError):
        return 0.0
    r, g, b = ((value >> 16) & 0xFF) / 255.0, ((value >> 8) & 0xFF) / 255.0, (value & 0xFF) / 255.0
    mx, mn = max(r, g, b), min(r, g, b)
    if mx == 0:
        return 0.0
    return (mx - mn) / mx


def _is_near_white(rgb: RGBColor, threshold: float = 0.94) -> bool:
    return _luminance(rgb) >= threshold


def _is_near_black(rgb: RGBColor, threshold: float = 0.12) -> bool:
    return _luminance(rgb) <= threshold


def _contrast_ratio(a: RGBColor, b: RGBColor) -> float:
    l1, l2 = _luminance(a), _luminance(b)
    lighter, darker = max(l1, l2), min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _quantize(rgb: RGBColor) -> str:
    try:
        value = _color_to_int(rgb)
    except (TypeError, ValueError):
        return ""
    r, g, b = (value >> 16) & 0xFF, (value >> 8) & 0xFF, value & 0xFF
    r = (r // 16) * 16 + 8
    g = (g // 16) * 16 + 8
    b = (b // 16) * 16 + 8
    return f"#{r:02X}{g:02X}{b:02X}"


def _hex_to_rgb(hex_value: str) -> RGBColor:
    cleaned = hex_value.strip().lstrip("#")
    return RGBColor(int(cleaned[0:2], 16), int(cleaned[2:4], 16), int(cleaned[4:6], 16))


def _theme_palette_from_pptx(template_bytes: bytes) -> List[RGBColor]:
    colors: List[RGBColor] = []
    with zipfile.ZipFile(io.BytesIO(template_bytes)) as archive:
        theme_names = sorted(
            n for n in archive.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml")
        )
        if not theme_names:
            return colors
        root = _parse_xml(archive.read(theme_names[0]))
        scheme = root.find(".//a:clrScheme", _NS)
        if scheme is None:
            return colors
        for tag in (
            "accent1",
            "accent2",
            "accent3",
            "accent4",
            "accent5",
            "accent6",
            "dk1",
            "dk2",
            "lt1",
            "lt2",
            "hlink",
        ):
            parsed = _parse_rgb(scheme.find(f"a:{tag}", _NS))
            if parsed is not None:
                colors.append(parsed)
    return colors


def _pptx_shape_colors(template_bytes: bytes, *, max_slides: int = 10) -> List[RGBColor]:
    prs = Presentation(io.BytesIO(template_bytes))
    counter: Counter[str] = Counter()
    for index in range(min(len(prs.slides), max_slides)):
        slide = prs.slides[index]
        bg = _slide_background_hex(slide)
        if bg:
            counter[_quantize(_hex_to_rgb(bg))] += 3
        for shape in slide.shapes:
            fill_hex = _shape_fill_hex(shape)
            if fill_hex:
                counter[_quantize(_hex_to_rgb(fill_hex))] += 2
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        try:
                            if run.font.color and run.font.color.rgb:
                                counter[_quantize(run.font.color.rgb)] += 1
                        except AttributeError:
                            pass
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.LINE and shape.line.color.rgb:
                    counter[_quantize(shape.line.color.rgb)] += 1
            except AttributeError:
                pass
    return [_hex_to_rgb(hex_key) for hex_key, _ in counter.most_common(24) if hex_key]


def _pdf_page_colors(template_bytes: bytes, *, max_pages: int = 3) -> List[RGBColor]:
    doc = fitz.open(stream=template_bytes, filetype="pdf")
    colors: List[RGBColor] = []
    try:
        for index in range(min(len(doc), max_pages)):
            accent, dark = _palette_from_page(doc[index])
            if accent is not None:
                colors.append(accent)
            if dark is not None:
                colors.append(dark)
            pixmap = doc[index].get_pixmap(matrix=fitz.Matrix(1.0, 1.0), alpha=False)
            image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
            for pixel, count in Counter(image.getdata()).most_common(32):
                if sum(pixel) >= 720 or sum(pixel) <= 30:
                    continue
                colors.append(RGBColor(*pixel))
    finally:
        doc.close()
    return colors


def extract_template_palette(
    template_bytes: bytes,
    template_file_type: str,
    *,
    max_colors: int = 6,
) -> List[str]:
    """Основные цвета шаблона (HEX), отсортированы по значимости."""
    raw: List[RGBColor] = []
    if template_file_type == "pptx":
        raw.extend(_theme_palette_from_pptx(template_bytes))
        raw.extend(_pptx_shape_colors(template_bytes))
    else:
        raw.extend(_pdf_page_colors(template_bytes))

    counter: Counter[str] = Counter()
    for rgb in raw:
        key = _quantize(rgb)
        if not key:
            continue
        if _is_near_white(_hex_to_rgb(key)) and len(counter) > 0:
            counter[key] += 1
        elif not _is_near_black(_hex_to_rgb(key)):
            counter[key] += 2
        else:
            counter[key] += 1

    ranked = [hex_key for hex_key, _ in counter.most_common(max_colors * 3)]
    result: List[str] = []
    for hex_key in ranked:
        rgb = _hex_to_rgb(hex_key)
        if _is_near_white(rgb) and len(result) >= 2:
            continue
        if hex_key not in result:
            result.append(hex_key)
        if len(result) >= max_colors:
            break

    if not result and ranked:
        result = ranked[:max_colors]

    key_accent = _pick_key_contrast_color(raw, ranked)
    if key_accent and key_accent not in result:
        if len(result) >= max_colors:
            result[-1] = key_accent
        else:
            result.append(key_accent)

    logger.info("Палитра шаблона (%s): %s", template_file_type, ", ".join(result))
    return result


def _pick_key_contrast_color(raw: List[RGBColor], ranked: List[str]) -> Optional[str]:
    if not raw:
        return None
    ranked_rgbs = [_hex_to_rgb(h) for h in ranked if h]
    if not ranked_rgbs:
        return None
    darkest = min(ranked_rgbs, key=_luminance)
    lightest = max(ranked_rgbs, key=_luminance)

    candidates: List[Tuple[float, str]] = []
    seen = set()
    for rgb in raw:
        key = _quantize(rgb)
        if not key or key in seen:
            continue
        seen.add(key)
        sample = _hex_to_rgb(key)
        sat = _saturation(sample)
        if sat < 0.52:
            continue
        contrast = max(_contrast_ratio(sample, darkest), _contrast_ratio(sample, lightest))
        score = sat * 1.4 + contrast
        candidates.append((score, key))

    if not candidates:
        return None
    candidates.sort(reverse=True)
    return candidates[0][1]


def apply_palette_to_style(style: UserTemplateStyle, palette_hex: List[str]) -> UserTemplateStyle:
    if not palette_hex:
        return style

    rgbs = [_hex_to_rgb(h) for h in palette_hex]
    non_white = [c for c in rgbs if not _is_near_white(c)]
    if not non_white:
        non_white = rgbs

    accent = max(non_white, key=_saturation)
    dark_candidates = [c for c in non_white if not _is_near_white(c)]
    dark = min(dark_candidates or non_white, key=_luminance)
    light = RGBColor(0xFF, 0xFF, 0xFF)

    card_candidates = [
        c for c in rgbs if 0.55 < _luminance(c) < 0.92 and _saturation(c) < 0.35
    ]
    card_fill = card_candidates[0] if card_candidates else light
    white_in_palette = next((c for c in rgbs if _is_near_white(c)), None)
    if white_in_palette is not None:
        card_fill = white_in_palette if _luminance(card_fill) > 0.94 else card_fill

    style.palette_hex = palette_hex
    style.accent_rgb = accent
    style.dark_rgb = dark
    style.light_rgb = RGBColor(0xFF, 0xFF, 0xFF)
    style.card_fill_rgb = card_fill
    if style.title_text_style.rgb is None:
        style.title_text_style = TextStyleHint(
            font_size_pt=style.title_text_style.font_size_pt,
            rgb=accent,
            bold=style.title_text_style.bold,
        )
    if style.body_text_style.rgb is None:
        style.body_text_style = TextStyleHint(
            font_size_pt=style.body_text_style.font_size_pt,
            rgb=dark,
            bold=style.body_text_style.bold,
        )
    return style

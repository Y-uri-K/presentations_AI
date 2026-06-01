from __future__ import annotations

import io
import logging
from collections import Counter
from dataclasses import dataclass, field
from typing import List, Optional

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.template_style_extractor import _NS, _parse_rgb

logger = logging.getLogger(__name__)

_MAX_SLIDES = 10


@dataclass
class SlideStyleSample:
    index: int
    background_hex: Optional[str] = None
    accent_colors: List[str] = field(default_factory=list)
    shape_kinds: List[str] = field(default_factory=list)
    has_background_image: bool = False
    title_size_pt: Optional[float] = None
    body_size_pt: Optional[float] = None


def _rgb_to_hex(rgb: Optional[RGBColor]) -> Optional[str]:
    if rgb is None:
        return None
    value = int(rgb)
    return f"#{value:06X}"


def _shape_kind_name(shape) -> str:
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "picture"
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return "line"
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return "auto_shape"
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            return "textbox"
        return str(shape.shape_type)
    except Exception:
        return "shape"


def _shape_fill_hex(shape) -> Optional[str]:
    try:
        fill = shape.fill
        if fill.type is not None and str(fill.type).endswith("SOLID"):
            return _rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def _slide_background_hex(slide) -> Optional[str]:
    try:
        c_sld = slide.element.find("p:cSld", _NS)
        if c_sld is None:
            return None
        bg = c_sld.find("p:bg", _NS)
        if bg is None:
            return None
        solid = bg.find(".//a:solidFill", _NS)
        return _rgb_to_hex(_parse_rgb(solid.find("a:srgbClr", _NS) if solid is not None else None))
    except Exception:
        return None


def inspect_pptx_slides(template_bytes: bytes, *, max_slides: int = _MAX_SLIDES) -> List[SlideStyleSample]:
    prs = Presentation(io.BytesIO(template_bytes))
    samples: List[SlideStyleSample] = []
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide_area = slide_width * slide_height if slide_width and slide_height else 0

    slide_count = min(len(prs.slides), max_slides)
    for index in range(slide_count):
        slide = prs.slides[index]
        colors: List[str] = []
        kinds: List[str] = []
        has_bg_image = False
        title_size: Optional[float] = None
        body_size: Optional[float] = None

        for shape in slide.shapes:
            kinds.append(_shape_kind_name(shape))
            color_hex = _shape_fill_hex(shape)
            if color_hex:
                colors.append(color_hex)
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and slide_area:
                    area = shape.width * shape.height
                    if area > slide_area * 0.45:
                        has_bg_image = True
            except Exception:
                pass
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is None:
                            continue
                        pt = run.font.size.pt
                        if title_size is None or pt > title_size:
                            title_size = pt
                        if body_size is None or pt < title_size:
                            body_size = pt

        samples.append(
            SlideStyleSample(
                index=index + 1,
                background_hex=_slide_background_hex(slide),
                accent_colors=sorted(set(colors))[:8],
                shape_kinds=sorted(set(kinds)),
                has_background_image=has_bg_image,
                title_size_pt=title_size,
                body_size_pt=body_size,
            )
        )
    return samples


def inspect_pdf_slides(template_bytes: bytes, *, max_slides: int = _MAX_SLIDES) -> List[SlideStyleSample]:
    doc = fitz.open(stream=template_bytes, filetype="pdf")
    samples: List[SlideStyleSample] = []
    try:
        for index in range(min(len(doc), max_slides)):
            page = doc[index]
            samples.append(
                SlideStyleSample(
                    index=index + 1,
                    background_hex="#FFFFFF",
                    accent_colors=[],
                    shape_kinds=["pdf_page"],
                    has_background_image=bool(page.get_images()),
                )
            )
    finally:
        doc.close()
    return samples


def inspect_template_slides(
    template_bytes: bytes,
    template_file_type: str,
    *,
    max_slides: int = _MAX_SLIDES,
) -> List[SlideStyleSample]:
    if template_file_type == "pdf":
        return inspect_pdf_slides(template_bytes, max_slides=max_slides)
    return inspect_pptx_slides(template_bytes, max_slides=max_slides)


def dominant_colors_from_samples(samples: List[SlideStyleSample]) -> tuple[Optional[str], Optional[str], Optional[str]]:
    counter: Counter[str] = Counter()
    for sample in samples:
        for color in sample.accent_colors:
            counter[color] += 1
        if sample.background_hex:
            counter[sample.background_hex] += 1
    if not counter:
        return None, None, None
    ranked = counter.most_common()
    accent = ranked[0][0] if ranked else None
    dark = ranked[1][0] if len(ranked) > 1 else accent
    light = ranked[2][0] if len(ranked) > 2 else None
    return accent, dark, light


def extract_largest_background_image_pptx(template_bytes: bytes, *, max_slides: int = _MAX_SLIDES) -> Optional[bytes]:
    """Извлекает крупнейшее изображение с фона среди первых слайдов PPTX."""
    prs = Presentation(io.BytesIO(template_bytes))
    best: tuple[int, bytes] | None = None
    slide_area = prs.slide_width * prs.slide_height

    slide_count = min(len(prs.slides), max_slides)
    for slide_index in range(slide_count):
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            try:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                area = shape.width * shape.height
                if slide_area and area < slide_area * 0.35:
                    continue
                blob = shape.image.blob
                if best is None or area > best[0]:
                    best = (area, blob)
            except Exception:
                continue
    return best[1] if best else None

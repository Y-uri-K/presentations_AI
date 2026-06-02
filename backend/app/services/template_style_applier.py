from __future__ import annotations

import io
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Emu, Pt

from app.services.pptx_fonts import (
    DEFAULT_PRESENTATION_FONT,
    apply_font_family_preserve_size,
    apply_font_family_preserve_to_slide,
)
from app.services.slide_palette_colors import SLIDE_BACKGROUND_WHITE, palette_rgb
from app.services.template_style_extractor import TextStyleHint, UserTemplateStyle


def apply_slide_background(
    slide,
    *,
    user_style: UserTemplateStyle,
    slide_width: int,
    slide_height: int,
) -> None:
    """Фон слайда — до контента и картинок."""
    if user_style.background_image:
        _apply_background_image(
            slide,
            user_style.background_image,
            slide_width=slide_width,
            slide_height=slide_height,
        )
    else:
        _apply_subtle_background(slide, user_style, slide_width=slide_width, slide_height=slide_height)


def apply_slide_typography(
    slide,
    *,
    user_style: UserTemplateStyle,
) -> None:
    _apply_style_to_all_text(slide, user_style)
    font_name = user_style.font_family or DEFAULT_PRESENTATION_FONT
    apply_font_family_preserve_to_slide(slide, font_name)


def apply_user_style_to_slide(
    slide,
    *,
    kind: str,
    user_style: UserTemplateStyle,
    slide_width: int,
    slide_height: int,
) -> None:
    apply_slide_background(
        slide, user_style=user_style, slide_width=slide_width, slide_height=slide_height
    )
    apply_slide_typography(slide, user_style=user_style)


def _apply_background_image(
    slide,
    image_bytes: bytes,
    *,
    slide_width: int,
    slide_height: int,
) -> None:
    from app.services.image_normalizer import normalize_image_for_pptx

    safe = normalize_image_for_pptx(image_bytes) or image_bytes
    stream = io.BytesIO(safe)
    stream.seek(0)
    slide.shapes.add_picture(stream, Emu(0), Emu(0), width=slide_width, height=slide_height)


def _apply_subtle_background(
    slide,
    user_style: UserTemplateStyle,
    *,
    slide_width: int,
    slide_height: int,
) -> None:
    fill_rgb = palette_rgb(user_style, "background")

    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(0),
        Emu(0),
        slide_width,
        slide_height,
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_rgb
    rect.line.fill.background()


def _color_to_int(color: RGBColor) -> int:
    try:
        return int(color)
    except TypeError:
        return (int(color[0]) << 16) | (int(color[1]) << 8) | int(color[2])


def _tint(rgb: RGBColor, amount: float) -> RGBColor:
    value = _color_to_int(rgb)
    r = (value >> 16) & 0xFF
    g = (value >> 8) & 0xFF
    b = value & 0xFF
    return RGBColor(
        min(255, int(r + (255 - r) * amount)),
        min(255, int(g + (255 - g) * amount)),
        min(255, int(b + (255 - b) * amount)),
    )


def _dominant_font_size_pt(text_frame) -> float:
    sizes = []
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
    return max(sizes) if sizes else 0.0


def _apply_style_to_all_text(slide, user_style: UserTemplateStyle) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            placeholder_type = shape.placeholder_format.type
        except (AttributeError, ValueError):
            placeholder_type = None

        is_title = placeholder_type == PP_PLACEHOLDER.TITLE or _dominant_font_size_pt(shape.text_frame) >= 24.0
        if is_title:
            hint = user_style.title_text_style
            fallback = palette_rgb(user_style, "accent")
            default_size = 32.0
        else:
            hint = user_style.body_text_style
            fallback = palette_rgb(user_style, "dark")
            default_size = 18.0

        _style_text_frame(
            shape.text_frame,
            hint,
            fallback_rgb=fallback,
            default_size_pt=default_size,
            preserve_existing_size=is_title and _dominant_font_size_pt(shape.text_frame) >= 24.0,
        )


def _style_text_frame(
    text_frame,
    hint: TextStyleHint,
    *,
    fallback_rgb: Optional[RGBColor],
    default_size_pt: float,
    preserve_existing_size: bool = False,
) -> None:
    """
    Только шрифт и цвет. Не вызывать configure_text_frame_wrap — он ставит
    MSO_AUTO_SIZE.NONE и обрезает текст в карточках после SHAPE_TO_FIT_TEXT.
    """
    autosize = getattr(text_frame, "auto_size", None)
    word_wrap = getattr(text_frame, "word_wrap", True)

    size_pt = hint.font_size_pt or default_size_pt
    color = hint.rgb or fallback_rgb

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text.strip():
                continue
            apply_font_family_preserve_size(run)
            if run.font.size is not None:
                pass
            elif not preserve_existing_size:
                run.font.size = Pt(size_pt)
            if hint.bold is not None:
                run.font.bold = hint.bold
            if color is not None:
                try:
                    existing = run.font.color.rgb
                    if existing is not None:
                        continue
                    run.font.color.rgb = color
                except AttributeError:
                    pass

    text_frame.word_wrap = word_wrap
    if autosize is not None:
        text_frame.auto_size = autosize

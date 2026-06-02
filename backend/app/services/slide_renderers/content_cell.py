from __future__ import annotations

"""Единый примитив ячейки: фон + textbox на всю выделенную область, без обрезки."""

from typing import Optional

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.services.pptx_fonts import apply_font_family_preserve_to_text_frame
from app.services.slide_palette_colors import palette_rgb
from app.services.slide_content_density import normalize_content_text
from app.services.slide_renderers.context import RenderContext
from app.services.slide_renderers.text_frame import configure_text_frame_wrap


def _rect_emu(
    ctx: RenderContext,
    left_pct: float,
    top_pct: float,
    width_pct: float,
    height_pct: float,
) -> tuple[int, int, int, int]:
    left = int(ctx.slide_width * left_pct)
    top = int(ctx.slide_height * top_pct)
    width = max(int(ctx.slide_width * width_pct), int(Emu(Inches(0.35))))
    height = max(int(ctx.slide_height * height_pct), int(Emu(Inches(0.28))))
    return left, top, width, height


def _scale_body_font_pt(lines: list[str], base_pt: float) -> float:
    total = sum(len(line) for line in lines)
    if total > 220:
        return max(8.0, base_pt - 3)
    if total > 140:
        return max(8.5, base_pt - 2)
    if total > 80:
        return max(9.0, base_pt - 1)
    return base_pt


def _write_paragraphs(
    text_frame,
    lines: list[str],
    *,
    font_pt: float,
    color,
    bold_first: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    configure_text_frame_wrap(text_frame, anchor=MSO_ANCHOR.TOP)
    text_frame.margin_left = Pt(12)
    text_frame.margin_right = Pt(12)
    text_frame.margin_top = Pt(10)
    text_frame.margin_bottom = Pt(10)

    visible = [normalize_content_text(line) for line in lines]
    visible = [line for line in visible if line]
    if not visible:
        visible = ["—"]

    for index, line in enumerate(visible):
        p = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_pt)
        run.font.color.rgb = color
        if bold_first and index == 0:
            run.font.bold = True
    apply_font_family_preserve_to_text_frame(text_frame)


def render_content_cell(
    ctx: RenderContext,
    *,
    left_pct: float,
    top_pct: float,
    width_pct: float,
    height_pct: float,
    body: str,
    heading: Optional[str] = None,
    kpi_value: Optional[str] = None,
    accent: bool = False,
    rounded: bool = True,
) -> None:
    """
    Ячейка контента: заливка на всю выделенную высоту, текст в textbox(ах) с переносом.
    kpi_value — крупное число сверху (~30%% высоты), body — остальное.
    """
    body = normalize_content_text(body) or ""
    if not body and not kpi_value:
        return

    left, top, width, height = _rect_emu(ctx, left_pct, top_pct, width_pct, height_pct)
    pad = Emu(Inches(0.04))

    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    bg = ctx.slide.shapes.add_shape(shape_type, left, top, width, height)
    bg.fill.solid()
    if accent:
        bg.fill.fore_color.rgb = ctx.accent
        text_color = palette_rgb(ctx.user_style, "on_accent")
        heading_color = palette_rgb(ctx.user_style, "on_accent")
    else:
        bg.fill.fore_color.rgb = palette_rgb(ctx.user_style, "card")
        text_color = palette_rgb(ctx.user_style, "muted")
        heading_color = palette_rgb(ctx.user_style, "dark")
    bg.line.color.rgb = ctx.accent
    try:
        bg.text_frame.clear()
    except Exception:
        pass

    inner_left = int(left) + pad
    inner_top = int(top) + pad
    inner_w = max(int(width) - 2 * int(pad), int(Emu(Inches(0.3))))
    inner_h = max(int(height) - 2 * int(pad), int(Emu(Inches(0.25))))

    body_lines: list[str] = []
    head = normalize_content_text(heading or "") if heading else ""
    if head and head not in ("—", "-", "•") and head != body:
        body_lines.append(head)
    for part in body.replace("\r", "").split("\n"):
        part = part.strip()
        if part:
            body_lines.append(part)
    if not body_lines:
        body_lines = [body or head or "—"]

    if kpi_value:
        value_h = max(int(inner_h * 0.32), int(Emu(Inches(0.45))))
        body_h = max(inner_h - value_h - int(Emu(Inches(0.04))), int(Emu(Inches(0.2))))
        val_box = ctx.slide.shapes.add_textbox(inner_left, inner_top, inner_w, value_h)
        _write_paragraphs(
            val_box.text_frame,
            [normalize_content_text(kpi_value) or "—"],
            font_pt=26 if accent else 24,
            color=heading_color if accent else ctx.accent,
            bold_first=True,
            align=PP_ALIGN.CENTER,
        )
        txt_box = ctx.slide.shapes.add_textbox(
            inner_left,
            inner_top + value_h + int(Emu(Inches(0.04))),
            inner_w,
            body_h,
        )
        body_pt = _scale_body_font_pt(body_lines, 11 if accent else 10)
        _write_paragraphs(
            txt_box.text_frame,
            body_lines,
            font_pt=body_pt,
            color=text_color,
        )
    else:
        txt_box = ctx.slide.shapes.add_textbox(inner_left, inner_top, inner_w, inner_h)
        body_pt = _scale_body_font_pt(body_lines, 12 if accent else 11)
        _write_paragraphs(
            txt_box.text_frame,
            body_lines,
            font_pt=body_pt,
            color=text_color,
            bold_first=bool(head and body_lines and body_lines[0] == head),
        )


def render_equal_cells_row(
    ctx: RenderContext,
    *,
    left_pct: float,
    top_pct: float,
    width_pct: float,
    height_pct: float,
    items: list[tuple[str, str, Optional[str]]],
    accent_index: int = 0,
) -> None:
    """Горизонтальный ряд ячеек одинаковой высоты (KPI / карточки)."""
    if not items:
        return
    gap = 0.014
    count = len(items)
    cell_w = (width_pct - gap * (count - 1)) / count
    for index, (heading, body, kpi) in enumerate(items):
        render_content_cell(
            ctx,
            left_pct=left_pct + index * (cell_w + gap),
            top_pct=top_pct,
            width_pct=cell_w,
            height_pct=height_pct,
            heading=heading,
            body=body,
            kpi_value=kpi,
            accent=index == accent_index,
            rounded=index % 2 == 0,
        )


def render_equal_cells_grid(
    ctx: RenderContext,
    bounds,
    items: list[tuple[str, str]],
    *,
    accent_index: int = 0,
    max_items: int = 4,
) -> None:
    """Сетка 2×2 (или 1×N), ячейки на всю высоту зоны без пустот."""
    visible = items[:max_items]
    if not visible:
        return
    count = len(visible)
    cols = 2 if count > 1 and not ctx.has_image_zone else 1
    rows = (count + cols - 1) // cols
    gap = 0.014
    span_w = bounds.right_pct - bounds.left_pct
    span_h = bounds.bottom_pct - bounds.top_pct
    cell_w = (span_w - gap * (cols - 1)) / cols
    cell_h = (span_h - gap * (rows - 1)) / rows

    for index, (heading, body) in enumerate(visible):
        row = index // cols
        col = index % cols
        render_content_cell(
            ctx,
            left_pct=bounds.left_pct + col * (cell_w + gap),
            top_pct=bounds.top_pct + row * (cell_h + gap),
            width_pct=cell_w,
            height_pct=cell_h,
            heading=heading,
            body=body,
            accent=index == accent_index,
            rounded=True,
        )


def render_equal_cells_stack(
    ctx: RenderContext,
    bounds,
    items: list[tuple[str, str]],
    *,
    accent_index: int = 0,
    max_items: int = 4,
) -> None:
    """Вертикальный стек: каждая строка получает равную долю высоты."""
    visible = [(normalize_content_text(h), normalize_content_text(b)) for h, b in items]
    visible = [(h, b) for h, b in visible if b and b not in ("—", "-")]
    visible = visible[:max_items]
    if not visible:
        return
    gap = 0.012
    span_h = bounds.bottom_pct - bounds.top_pct
    cell_h = (span_h - gap * (len(visible) - 1)) / len(visible)
    span_w = bounds.right_pct - bounds.left_pct

    for index, (heading, body) in enumerate(visible):
        render_content_cell(
            ctx,
            left_pct=bounds.left_pct,
            top_pct=bounds.top_pct + index * (cell_h + gap),
            width_pct=span_w,
            height_pct=cell_h,
            heading=heading,
            body=body,
            accent=index == accent_index,
            rounded=index % 2 == 0,
        )

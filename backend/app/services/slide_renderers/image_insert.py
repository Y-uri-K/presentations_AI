from __future__ import annotations

import io
import logging

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from app.services.image_normalizer import normalize_image_for_pptx
from app.services.slide_renderers.context import RenderContext
from app.services.slide_renderers.drawing import _emu_fraction, _fit_picture_size
from app.services.slide_renderers.layout_bounds import image_bounds_for_slide

logger = logging.getLogger(__name__)

_MIN_PICTURE_EMU = Emu(91440)


def count_pictures_on_slide(slide) -> int:
    count = 0
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                count += 1
        except (AttributeError, ValueError):
            continue
    return count


def _insert_picture_bytes(
    slide,
    *,
    image_bytes: bytes,
    slide_width: int,
    slide_height: int,
    left_pct: float,
    width_pct: float,
    top_pct: float,
    bottom_pct: float,
) -> bool:
    jpeg_bytes = normalize_image_for_pptx(image_bytes)
    if not jpeg_bytes:
        return False

    left, box_width = _emu_fraction(slide_width, left_pct, left_pct + width_pct)
    top, box_height = _emu_fraction(slide_height, top_pct, bottom_pct)
    if int(box_width) <= int(_MIN_PICTURE_EMU) or int(box_height) <= int(_MIN_PICTURE_EMU):
        return False

    pic_width, pic_height = _fit_picture_size(jpeg_bytes, int(box_width), int(box_height))
    pic_width = int(pic_width)
    pic_height = int(pic_height)
    if pic_width < int(_MIN_PICTURE_EMU) or pic_height < int(_MIN_PICTURE_EMU):
        return False

    left = int(left) + max(0, (int(box_width) - pic_width) // 2)
    top = int(top) + max(0, (int(box_height) - pic_height) // 2)

    slide_w = int(slide_width)
    slide_h = int(slide_height)
    if left + pic_width > slide_w:
        left = max(0, slide_w - pic_width)
    if top + pic_height > slide_h:
        top = max(0, slide_h - pic_height)

    try:
        stream = io.BytesIO(jpeg_bytes)
        stream.seek(0)
        slide.shapes.add_picture(stream, left, top, width=pic_width, height=pic_height)
    except Exception as exc:
        logger.error("add_picture: %s", exc)
        return False

    return True


def insert_slide_image_early(ctx: RenderContext) -> bool:
    """
    Вставка сразу после фона, до контента.
    Без spTree.remove/insert — иначе PowerPoint показывает ошибку восстановления.
    """
    if not ctx.image_bytes:
        return False
    if ctx.spec.type == "table":
        return True

    bounds = image_bounds_for_slide(ctx)
    return _insert_picture_bytes(
        ctx.slide,
        image_bytes=ctx.image_bytes,
        slide_width=int(ctx.slide_width),
        slide_height=int(ctx.slide_height),
        left_pct=bounds.left_pct,
        width_pct=bounds.right_pct - bounds.left_pct,
        top_pct=bounds.top_pct,
        bottom_pct=bounds.bottom_pct,
    )


def insert_slide_image_verified(ctx: RenderContext) -> bool:
    """Повторная вставка только если ранней не было (fallback)."""
    if ctx.spec.type == "table":
        return True

    before = count_pictures_on_slide(ctx.slide)
    if before > 0:
        logger.debug("Слайд %s: картинка уже на слайде (%s)", ctx.slide_index + 1, before)
        return True

    bounds = image_bounds_for_slide(ctx)
    width_pct = bounds.right_pct - bounds.left_pct
    height_pct = bounds.bottom_pct - bounds.top_pct

    ok = _insert_picture_bytes(
        ctx.slide,
        image_bytes=ctx.image_bytes,
        slide_width=int(ctx.slide_width),
        slide_height=int(ctx.slide_height),
        left_pct=bounds.left_pct,
        width_pct=width_pct,
        top_pct=bounds.top_pct,
        bottom_pct=bounds.bottom_pct,
    )
    if ok and count_pictures_on_slide(ctx.slide) > before:
        logger.info("Слайд %s: картинка вставлена (fallback)", ctx.slide_index + 1)
        return True

    logger.error("Слайд %s: вставка изображения не удалась", ctx.slide_index + 1)
    return False

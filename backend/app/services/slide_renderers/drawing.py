from __future__ import annotations

import io
from typing import Iterable, List, Optional, Sequence, Tuple

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.services.pptx_fonts import apply_font_family_preserve_to_text_frame
from app.services.slide_palette_colors import palette_rgb
from app.services.slide_renderers.context import RenderContext
from app.services.slide_renderers.layout_bounds import image_bounds_for_slide
from app.services.slide_content_density import normalize_content_text
from app.services.slide_renderers.text_frame import (
    configure_text_frame_autofit_grow,
    configure_text_frame_wrap,
)
from app.services.slide_renderers.text_layout import (
    MIN_CARD_HEIGHT_PCT,
    card_content_height_pct,
    estimate_block_height_pct,
)


def _emu_fraction(total: int, start_pct: float, end_pct: float) -> Tuple[int, int]:
    start = int(total * start_pct)
    size = int(total * (end_pct - start_pct))
    return start, max(size, Emu(Inches(0.1)))


def add_brand_accent_bar(ctx: RenderContext) -> None:
    """Верхняя акцентная полоса в духе корпоративного шаблона."""
    if "accent_top_bar" not in (ctx.user_style.key_elements or []):
        return
    accent = ctx.accent
    height = max(int(ctx.slide_height * 0.018), Emu(Inches(0.08)))
    bar = ctx.slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(0),
        Emu(0),
        ctx.slide_width,
        height,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()


def add_title(ctx: RenderContext, *, top_pct: float = 0.06, height_pct: float = 0.16) -> None:
    left, width = _emu_fraction(ctx.slide_width, 0.06, 0.94)
    top, height = _emu_fraction(ctx.slide_height, top_pct, top_pct + height_pct)
    box = ctx.slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    configure_text_frame_wrap(frame, anchor=MSO_ANCHOR.TOP)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = ctx.spec.title
    run.font.bold = True
    size_pt = ctx.user_style.title_text_style.font_size_pt or 32
    run.font.size = Pt(size_pt)
    run.font.color.rgb = ctx.title_color
    apply_font_family_preserve_to_text_frame(frame)


def add_subtitle(ctx: RenderContext, text: str, *, top_pct: float = 0.2) -> None:
    if not text:
        return
    left, width = _emu_fraction(ctx.slide_width, 0.06, 0.94)
    top, height = _emu_fraction(ctx.slide_height, top_pct, top_pct + 0.14)
    box = ctx.slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    configure_text_frame_wrap(frame, anchor=MSO_ANCHOR.TOP)
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(ctx.user_style.body_text_style.font_size_pt or 20)
    run.font.color.rgb = palette_rgb(ctx.user_style, "muted")
    apply_font_family_preserve_to_text_frame(frame)


def _fill_card_text(
    text_frame,
    *,
    heading: str,
    body: str,
    accent: bool,
    ctx: RenderContext,
) -> None:
    configure_text_frame_wrap(text_frame, anchor=MSO_ANCHOR.TOP)
    text_frame.margin_left = Pt(10)
    text_frame.margin_right = Pt(10)
    text_frame.margin_top = Pt(8)
    text_frame.margin_bottom = Pt(8)

    show_heading = (
        heading
        and heading not in ("—", "-", "•")
        and heading != body
        and len(heading) <= 50
        and len(heading.split()) <= 10
    )
    p1 = text_frame.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    if show_heading:
        r1 = p1.add_run()
        r1.text = heading
        r1.font.bold = True
        r1.font.size = Pt(13 if not accent else 15)
        r1.font.color.rgb = (
            palette_rgb(ctx.user_style, "on_accent") if accent else palette_rgb(ctx.user_style, "dark")
        )
        paragraphs = body.split("\n")
        for line_index, line in enumerate(paragraphs):
            line = line.strip()
            if not line:
                continue
            p_body = text_frame.add_paragraph()
            p_body.alignment = PP_ALIGN.LEFT
            r2 = p_body.add_run()
            r2.text = line
            r2.font.size = Pt(11 if not accent else 12)
            r2.font.color.rgb = (
                palette_rgb(ctx.user_style, "on_accent") if accent else palette_rgb(ctx.user_style, "muted")
            )
    else:
        paragraphs = [p.strip() for p in body.split("\n") if p.strip()] or [body]
        for line_index, line in enumerate(paragraphs):
            p = p1 if line_index == 0 else text_frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = line
            r.font.size = Pt(11 if not accent else 12)
            r.font.color.rgb = (
                palette_rgb(ctx.user_style, "on_accent") if accent else palette_rgb(ctx.user_style, "muted")
            )
    apply_font_family_preserve_to_text_frame(text_frame)


def _clear_text_frame(text_frame) -> None:
    """Пустая фигура без видимого текста."""
    text_frame.clear()
    text_frame.paragraphs[0].text = ""


def add_card(
    ctx: RenderContext,
    *,
    left_pct: float,
    top_pct: float,
    width_pct: float,
    height_pct: float,
    heading: str,
    body: str,
    accent: bool = False,
    style: str = "rounded",
) -> None:
    """Карточка на всю выделенную область (делегат content_cell)."""
    del style
    from app.services.slide_renderers.content_cell import render_content_cell

    heading = normalize_content_text(heading) or ""
    body = normalize_content_text(body) or heading
    if not body or body in ("—", "-"):
        return
    render_content_cell(
        ctx,
        left_pct=left_pct,
        top_pct=top_pct,
        width_pct=width_pct,
        height_pct=height_pct,
        heading=heading,
        body=body,
        accent=accent,
        rounded=True,
    )


def add_card_stack(
    ctx: RenderContext,
    *,
    left_pct: float,
    top_pct: float,
    width_pct: float,
    bottom_pct: float,
    items: Sequence[Tuple[str, str]],
    accent_first: bool = True,
    accent_index: Optional[int] = None,
    style: str = "sidebar",
    gap_pct: float = 0.014,
) -> None:
    """Вертикальный стек: равные доли высоты, полный текст в каждой ячейке."""
    del style, gap_pct
    from app.services.slide_renderers.content_cell import render_equal_cells_stack
    from app.services.slide_renderers.layout_bounds import ContentBounds

    visible = [(normalize_content_text(h), normalize_content_text(b)) for h, b in items]
    visible = [(h, b) for h, b in visible if b and b not in ("—", "-")]
    if not visible:
        return
    if accent_index is None:
        accent_index = 0 if accent_first else -1
    bounds = ContentBounds(
        left_pct=left_pct,
        right_pct=left_pct + width_pct,
        top_pct=top_pct,
        bottom_pct=bottom_pct,
    )
    render_equal_cells_stack(ctx, bounds, visible, accent_index=accent_index)


def add_stacked_text_block(
    ctx: RenderContext,
    *,
    left_pct: float,
    top_pct: float,
    width_pct: float,
    height_pct: float,
    items: Sequence[Tuple[str, str]],
) -> None:
    """
    Один текстовый блок на всю колонку — перенос строк без обрезки по краям карточек.
    """
    if not items:
        return
    left, width = _emu_fraction(ctx.slide_width, left_pct, left_pct + width_pct)
    top, height = _emu_fraction(ctx.slide_height, top_pct, top_pct + height_pct)
    # Минимальная высота; SHAPE_TO_FIT_TEXT расширит блок под весь текст.
    min_h = Emu(Inches(0.4))
    box = ctx.slide.shapes.add_textbox(left, top, width, max(int(height), int(min_h)))
    frame = box.text_frame
    configure_text_frame_autofit_grow(frame, anchor=MSO_ANCHOR.TOP)
    frame.margin_left = Pt(8)
    frame.margin_right = Pt(8)
    frame.margin_top = Pt(6)
    frame.margin_bottom = Pt(6)

    for index, (heading, body) in enumerate(items):
        text = normalize_content_text(body) or normalize_content_text(heading)
        if not text:
            continue
        h = (heading or "").strip()
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.space_after = Pt(8)
        p.level = 0
        if h and h not in ("•", "—", "1", "2", "3", "4", "5") and h != text and len(h) < 50:
            r_head = p.add_run()
            r_head.text = f"{h}\n"
            r_head.font.bold = True
            r_head.font.size = Pt(11)
            r_head.font.color.rgb = palette_rgb(ctx.user_style, "dark")
        r_body = p.add_run()
        r_body.text = text
        r_body.font.size = Pt(10)
        r_body.font.color.rgb = palette_rgb(ctx.user_style, "muted")
    apply_font_family_preserve_to_text_frame(frame)


def add_metric_block(
    ctx: RenderContext,
    *,
    left_pct: float,
    top_pct: float,
    width_pct: float,
    height_pct: float,
    value: str,
    label: str,
    note: Optional[str] = None,
) -> None:
    from app.services.slide_renderers.content_cell import render_content_cell

    render_content_cell(
        ctx,
        left_pct=left_pct,
        top_pct=top_pct,
        width_pct=width_pct,
        height_pct=height_pct,
        heading=label,
        body=(note or "").strip(),
        kpi_value=value,
        accent=False,
        rounded=True,
    )


def add_horizontal_timeline(
    ctx: RenderContext,
    steps: Sequence[Tuple[str, Optional[str]]],
    *,
    top_pct: float = 0.38,
) -> None:
    count = len(steps)
    if count == 0:
        return
    y_line = int(ctx.slide_height * top_pct)
    x_start = int(ctx.slide_width * 0.08)
    x_end = int(ctx.slide_width * 0.92)
    segment = (x_end - x_start) // max(count - 1, 1)

    line = ctx.slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        x_start,
        y_line,
        x_end - x_start,
        Emu(Inches(0.03)),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ctx.accent
    line.line.fill.background()

    for index, (label, description) in enumerate(steps):
        x = x_start + segment * index
        node = ctx.slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            x - Emu(Inches(0.12)),
            y_line - Emu(Inches(0.12)),
            Emu(Inches(0.24)),
            Emu(Inches(0.24)),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = ctx.accent
        node.line.fill.background()

        box_top = y_line + Emu(Inches(0.2))
        box = ctx.slide.shapes.add_textbox(
            x - Emu(Inches(0.9)),
            box_top,
            Emu(Inches(1.8)),
            Emu(Inches(0.9)),
        )
        tf = box.text_frame
        configure_text_frame_wrap(tf, anchor=MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.bold = True
        r.font.size = Pt(11)
        r.font.color.rgb = palette_rgb(ctx.user_style, "dark")
        if description:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = description
            r2.font.size = Pt(9)
            r2.font.color.rgb = palette_rgb(ctx.user_style, "muted")
        apply_font_family_preserve_to_text_frame(tf)


def add_process_flow(ctx: RenderContext, steps: Sequence[Tuple[str, Optional[str]]]) -> None:
    from app.services.slide_renderers.content_cell import render_equal_cells_row
    from app.services.slide_renderers.layout_bounds import content_bounds_for_slide

    count = len(steps)
    if count == 0:
        return
    bounds = content_bounds_for_slide(ctx)
    items = [
        (title, (description or "").strip() or title, None)
        for title, description in steps[:4]
    ]
    render_equal_cells_row(
        ctx,
        left_pct=bounds.left_pct,
        top_pct=bounds.top_pct,
        width_pct=bounds.right_pct - bounds.left_pct,
        height_pct=bounds.bottom_pct - bounds.top_pct,
        items=items,
        accent_index=count - 1,
    )


def add_table(
    ctx: RenderContext,
    headers: List[str],
    rows: List[List[str]],
    *,
    top_pct: float = 0.32,
) -> None:
    if not headers:
        return
    cols = len(headers)
    table_rows = 1 + len(rows)
    left, width = _emu_fraction(ctx.slide_width, 0.06, 0.94)
    top, height = _emu_fraction(ctx.slide_height, top_pct, 0.9)
    table_shape = ctx.slide.shapes.add_table(table_rows, cols, left, top, width, height)
    table = table_shape.table

    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        configure_text_frame_wrap(cell.text_frame, anchor=MSO_ANCHOR.MIDDLE)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(11)
                run.font.color.rgb = palette_rgb(ctx.user_style, "on_accent")
        cell.fill.solid()
        cell.fill.fore_color.rgb = ctx.accent

    for row_index, row in enumerate(rows, start=1):
        for col_index in range(cols):
            value = row[col_index] if col_index < len(row) else ""
            cell = table.cell(row_index, col_index)
            cell.text = value
            configure_text_frame_wrap(cell.text_frame, anchor=MSO_ANCHOR.TOP)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = palette_rgb(ctx.user_style, "dark")


def add_logo_zone_marker(ctx: RenderContext) -> None:
    """Маркер зоны логотипа на титуле (по анализу шаблона)."""
    zone = ctx.user_style.title_logo_zone
    if not zone or ctx.spec.type != "title":
        return
    left_pct = 0.06 if zone == "top_left" else 0.78 if zone == "top_right" else 0.38
    left, width = _emu_fraction(ctx.slide_width, left_pct, left_pct + 0.14)
    top, height = _emu_fraction(ctx.slide_height, 0.05, 0.16)
    shape = ctx.slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ctx.accent
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = " "
    r.font.size = Pt(8)


def count_pictures_on_slide(slide) -> int:
    count = 0
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                count += 1
        except (AttributeError, ValueError):
            continue
    return count


def _fit_picture_size(
    image_bytes: bytes,
    box_width_emu: int,
    box_height_emu: int,
) -> tuple[int, int]:
    from PIL import Image

    with Image.open(io.BytesIO(image_bytes)) as img:
        img_w, img_h = img.size
    box_w = max(int(box_width_emu), 1)
    box_h = max(int(box_height_emu), 1)
    if img_w <= 0 or img_h <= 0:
        return box_w, box_h
    aspect = img_w / img_h
    box_aspect = box_w / box_h
    if aspect >= box_aspect:
        pic_w = box_w
        pic_h = max(1, int(box_w / aspect))
    else:
        pic_h = box_h
        pic_w = max(1, int(box_h * aspect))
    return pic_w, pic_h


def add_image_area(
    ctx: RenderContext,
    *,
    left_pct: float = 0.52,
    width_pct: float = 0.42,
    top_pct: float = 0.32,
    bottom_pct: float = 0.88,
    behind_content: bool = True,
) -> None:
    del behind_content
    from app.services.slide_renderers.image_insert import insert_slide_image_verified

    insert_slide_image_verified(ctx)


def add_slide_image(ctx: RenderContext) -> None:
    """Картинка в отдельной зоне — под текстом (z-order)."""
    add_slide_image_verified(ctx)


def add_slide_image_verified(ctx: RenderContext) -> bool:
    """Вставляет изображение после контента; возвращает True при успехе."""
    from app.services.slide_renderers.image_insert import insert_slide_image_verified

    return insert_slide_image_verified(ctx)


def grid_positions(count: int) -> Iterable[Tuple[float, float, float, float]]:
    layouts = {
        2: [(0.06, 0.34, 0.42, 0.5), (0.52, 0.34, 0.42, 0.5)],
        3: [
            (0.06, 0.34, 0.27, 0.5),
            (0.365, 0.34, 0.27, 0.5),
            (0.67, 0.34, 0.27, 0.5),
        ],
        4: [
            (0.06, 0.34, 0.42, 0.28),
            (0.52, 0.34, 0.42, 0.28),
            (0.06, 0.66, 0.42, 0.28),
            (0.52, 0.66, 0.42, 0.28),
        ],
    }
    if count in layouts:
        yield from layouts[count]
        return
    cols = 3
    rows = (count + cols - 1) // cols
    cell_w = 0.88 / cols
    cell_h = 0.55 / max(rows, 1)
    for index in range(count):
        row = index // cols
        col = index % cols
        yield (
            0.06 + col * cell_w,
            0.34 + row * cell_h,
            cell_w - 0.02,
            cell_h - 0.02,
        )

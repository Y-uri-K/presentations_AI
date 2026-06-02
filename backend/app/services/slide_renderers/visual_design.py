from __future__ import annotations

"""Визуальные примитивы: акценты, KPI-чипы, нумерация — плотная информационная подача."""

from typing import List, Sequence, Tuple

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.services.pptx_fonts import apply_font_family_preserve_to_text_frame
from app.services.slide_palette_colors import palette_rgb
from app.services.slide_renderers.context import RenderContext
from app.services.slide_renderers.layout_bounds import ContentBounds
from app.services.slide_renderers.text_frame import configure_text_frame_wrap


def _frac(total: int, start_pct: float, end_pct: float) -> tuple[int, int]:
    start = int(total * start_pct)
    size = int(total * (end_pct - start_pct))
    return start, max(size, Emu(Inches(0.08)))


def add_content_zone_frame(ctx: RenderContext, bounds: ContentBounds) -> None:
    """Тонкая рамка контентной зоны — визуальная «сетка» слайда."""
    left, width = _frac(ctx.slide_width, bounds.left_pct, bounds.right_pct)
    top, height = _frac(ctx.slide_height, bounds.top_pct, bounds.bottom_pct)
    frame = ctx.slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top,
        width,
        height,
    )
    frame.fill.background()
    frame.line.color.rgb = palette_rgb(ctx.user_style, "muted")
    frame.line.width = Pt(0.75)


def add_left_accent_rail(ctx: RenderContext, bounds: ContentBounds, *, width_pct: float = 0.008) -> None:
    left, w = _frac(ctx.slide_width, bounds.left_pct, bounds.left_pct + width_pct)
    top, height = _frac(ctx.slide_height, bounds.top_pct, bounds.bottom_pct)
    rail = ctx.slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, w, height)
    rail.fill.solid()
    rail.fill.fore_color.rgb = ctx.accent
    rail.line.fill.background()


def add_stat_chip_row(
    ctx: RenderContext,
    chips: Sequence[Tuple[str, str]],
    *,
    left_pct: float,
    top_pct: float,
    width_pct: float,
    height_pct: float = 0.11,
) -> None:
    """Горизонтальный ряд KPI-чипов (значение + подпись)."""
    visible = [(v.strip(), l.strip()) for v, l in chips if v.strip() or l.strip()][:4]
    if not visible:
        return
    gap = 0.015
    count = len(visible)
    chip_w = (width_pct - gap * (count - 1)) / count
    for index, (value, label) in enumerate(visible):
        left = left_pct + index * (chip_w + gap)
        l, w = _frac(ctx.slide_width, left, left + chip_w)
        t, h = _frac(ctx.slide_height, top_pct, top_pct + height_pct)
        shape = ctx.slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
        accent_chip = index == 0
        shape.fill.solid()
        shape.fill.fore_color.rgb = ctx.accent if accent_chip else palette_rgb(ctx.user_style, "card")
        shape.line.color.rgb = ctx.accent
        tf = shape.text_frame
        configure_text_frame_wrap(tf, anchor=MSO_ANCHOR.MIDDLE)
        tf.margin_left = Pt(6)
        tf.margin_right = Pt(6)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r_val = p.add_run()
        r_val.text = value
        r_val.font.bold = True
        r_val.font.size = Pt(16 if accent_chip else 14)
        r_val.font.color.rgb = (
            palette_rgb(ctx.user_style, "on_accent") if accent_chip else ctx.accent
        )
        if label:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r_lbl = p2.add_run()
            r_lbl.text = label
            r_lbl.font.size = Pt(9)
            r_lbl.font.color.rgb = (
                palette_rgb(ctx.user_style, "on_accent") if accent_chip else palette_rgb(ctx.user_style, "muted")
            )
        apply_font_family_preserve_to_text_frame(tf)


def add_number_badge(
    ctx: RenderContext,
    number: str,
    *,
    left_pct: float,
    top_pct: float,
    size_pct: float = 0.055,
) -> None:
    """Круглый номер пункта."""
    cx = int(ctx.slide_width * left_pct)
    cy = int(ctx.slide_height * top_pct)
    size = max(int(ctx.slide_height * size_pct), int(Emu(Inches(0.28))))
    shape = ctx.slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        cx,
        cy,
        size,
        size,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ctx.accent
    shape.line.fill.background()
    tf = shape.text_frame
    configure_text_frame_wrap(tf, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(number)
    r.font.bold = True
    r.font.size = Pt(12)
    r.font.color.rgb = palette_rgb(ctx.user_style, "on_accent")


def add_hub_center(
    ctx: RenderContext,
    label: str,
    *,
    center_x_pct: float,
    center_y_pct: float,
    radius_pct: float = 0.07,
) -> None:
    """Центральный узел hub-диаграммы."""
    cx = int(ctx.slide_width * center_x_pct)
    cy = int(ctx.slide_height * center_y_pct)
    r = max(int(min(ctx.slide_width, ctx.slide_height) * radius_pct), int(Emu(Inches(0.45))))
    shape = ctx.slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        cx - r // 2,
        cy - r // 2,
        r,
        r,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ctx.accent
    shape.line.fill.background()
    tf = shape.text_frame
    configure_text_frame_wrap(tf, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r_run = p.add_run()
    r_run.text = (label or "●")[:40]
    r_run.font.bold = True
    r_run.font.size = Pt(11)
    r_run.font.color.rgb = palette_rgb(ctx.user_style, "on_accent")
    apply_font_family_preserve_to_text_frame(tf)


def add_connector_line(
    ctx: RenderContext,
    x1_pct: float,
    y1_pct: float,
    x2_pct: float,
    y2_pct: float,
) -> None:
    x1, y1 = int(ctx.slide_width * x1_pct), int(ctx.slide_height * y1_pct)
    x2, y2 = int(ctx.slide_width * x2_pct), int(ctx.slide_height * y2_pct)
    left = min(x1, x2)
    top = min(y1, y2)
    width = max(abs(x2 - x1), int(Emu(Inches(0.02))))
    height = max(abs(y2 - y1), int(Emu(Inches(0.02))))
    line = ctx.slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top,
        width,
        height,
    )
    line.fill.solid()
    line.fill.fore_color.rgb = palette_rgb(ctx.user_style, "muted")
    line.line.fill.background()


def extract_metric_chips_from_cards(cards: list) -> List[Tuple[str, str]]:
    """Пытается вытащить KPI из highlight или title карточек."""
    import re

    metric_re = re.compile(r"^[\d+\-~≈]|[\d]+%|\d+\s*(млн|млрд|тыс|k|K|M)", re.IGNORECASE)
    chips: List[Tuple[str, str]] = []
    for card in cards[:4]:
        title = getattr(card, "title", "") or ""
        text = getattr(card, "text", "") or ""
        highlight = getattr(card, "highlight", None) or ""
        candidate = str(highlight or title).strip()
        if metric_re.search(candidate):
            chips.append((candidate[:20], (text or title)[:40]))
        elif metric_re.search(text[:30]):
            chips.append((title[:20] or "—", text[:40]))
    return chips

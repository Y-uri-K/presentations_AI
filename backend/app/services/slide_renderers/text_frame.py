from __future__ import annotations

from typing import Any

from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE


def _apply_body_pr(text_frame: Any, *, anchor: MSO_ANCHOR) -> None:
    try:
        body_pr = text_frame._txBody.bodyPr
        body_pr.set("wrap", "square")
        body_pr.set("anchor", "t" if anchor == MSO_ANCHOR.TOP else "ctr")
    except AttributeError:
        pass


def configure_text_frame_wrap(
    text_frame: Any,
    *,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    """Перенос по словам; фиксированный размер фигуры."""
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = anchor
    _apply_body_pr(text_frame, anchor=anchor)


def configure_text_frame_autofit_grow(
    text_frame: Any,
    *,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    """Текстовый блок подстраивает высоту под содержимое (без обрезки снизу)."""
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    text_frame.vertical_anchor = anchor
    _apply_body_pr(text_frame, anchor=anchor)

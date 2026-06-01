from __future__ import annotations

import io
import logging
from typing import Dict, List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from app.schemas.template_blueprint import BlueprintSlide, PresentationBlueprint
from app.services.image_normalizer import normalize_image_for_pptx
from app.services.template_driven.structure_analyzer import AnalyzedTemplate, TemplateSlideRuntime
from app.services.template_driven.template_slide_clone import clone_slide_from_template
from app.services.template_driven.text_utils import is_placeholder_text
from app.services.template_style_applier import apply_slide_typography

logger = logging.getLogger(__name__)


def _clear_slides(prs: Presentation) -> None:
    for index in range(len(prs.slides) - 1, -1, -1):
        slide_id = prs.slides._sldIdLst[index]
        prs.part.drop_rel(slide_id.rId)
        del prs.slides._sldIdLst[index]


def _set_placeholder_text(slide, ph_idx: int, text: str) -> bool:
    for shape in slide.placeholders:
        try:
            if shape.placeholder_format.idx == ph_idx:
                shape.text_frame.clear()
                shape.text = text
                return True
        except (AttributeError, ValueError):
            continue
    return False


def _fill_text_shape_by_index(slide, shape_index: int, text: str) -> None:
    extra = 0
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        if extra == shape_index:
            shape.text_frame.clear()
            shape.text = text
            return
        extra += 1


def _card_texts(slide: BlueprintSlide) -> List[str]:
    texts: List[str] = []
    for card in slide.cards:
        block = f"{card.title}\n{card.text}".strip() if card.text else card.title
        texts.append(block)
    return texts


def _fill_slide(slide, bp: BlueprintSlide, runtime: TemplateSlideRuntime) -> None:
    schema = runtime.schema

    if schema.title_slot and bp.title:
        _set_placeholder_text(slide, int(schema.title_slot.split("_")[1]), bp.title)
    else:
        try:
            if slide.shapes.title:
                slide.shapes.title.text = bp.title
        except AttributeError:
            pass

    if schema.subtitle_slot and bp.subtitle:
        _set_placeholder_text(slide, int(schema.subtitle_slot.split("_")[1]), bp.subtitle)
    elif bp.subtitle:
        for shape in slide.placeholders:
            try:
                if shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                    shape.text = bp.subtitle
                    break
            except (AttributeError, ValueError):
                continue

    for slot_id, text in bp.slot_texts.items():
        if not text or is_placeholder_text(text):
            continue
        if slot_id.startswith("ph_"):
            try:
                idx = int(slot_id.split("_", 1)[1])
                _set_placeholder_text(slide, idx, text)
            except ValueError:
                pass
        elif slot_id.startswith("shape_"):
            try:
                _fill_text_shape_by_index(slide, int(slot_id.split("_", 1)[1]), text)
            except ValueError:
                pass

    if bp.bullets:
        body_ph = None
        for shape in slide.placeholders:
            try:
                if shape.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                    body_ph = shape
                    break
            except (AttributeError, ValueError):
                continue
        if body_ph is not None:
            tf = body_ph.text_frame
            tf.clear()
            for index, bullet in enumerate(bp.bullets):
                p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
                p.text = bullet
                p.level = 0

    card_texts = _card_texts(bp)
    for card_index, slot_id in enumerate(schema.card_slots):
        if card_index < len(card_texts):
            text = card_texts[card_index]
            if slot_id.startswith("shape_"):
                _fill_text_shape_by_index(slide, int(slot_id.split("_", 1)[1]), text)
            elif slot_id.startswith("ph_"):
                _set_placeholder_text(slide, int(slot_id.split("_", 1)[1]), text)

    if card_texts and not schema.card_slots:
        for index, text in enumerate(card_texts):
            _fill_text_shape_by_index(slide, index, text)

    if bp.comparison and schema.column_slots:
        columns = [
            f"{bp.comparison.left_heading}\n" + "\n".join(f"• {p}" for p in bp.comparison.left_points),
            f"{bp.comparison.right_heading}\n" + "\n".join(f"• {p}" for p in bp.comparison.right_points),
        ]
        for slot_id, text in zip(schema.column_slots, columns):
            if slot_id.startswith("shape_"):
                _fill_text_shape_by_index(slide, int(slot_id.split("_", 1)[1]), text)

    if bp.table and schema.has_table:
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
                continue
            table = shape.table
            headers = bp.table.headers
            for col_index, header in enumerate(headers[: len(table.columns)]):
                table.cell(0, col_index).text = header
            for row_index, row in enumerate(bp.table.rows[: len(table.rows) - 1], start=1):
                for col_index in range(len(table.columns)):
                    value = row[col_index] if col_index < len(row) else ""
                    table.cell(row_index, col_index).text = value
            break

    if bp.speaker_notes:
        try:
            notes = slide.notes_slide
            if notes.notes_text_frame:
                notes.notes_text_frame.text = bp.speaker_notes
        except Exception:
            pass


def _insert_image(slide, image_bytes: bytes, slide_width: int, slide_height: int) -> None:
    jpeg = normalize_image_for_pptx(image_bytes)
    if not jpeg:
        return
    stream = io.BytesIO(jpeg)
    stream.seek(0)
    left = int(slide_width * 0.55)
    top = int(slide_height * 0.30)
    width = int(slide_width * 0.40)
    height = int(slide_height * 0.55)
    slide.shapes.add_picture(stream, left, top, width=width, height=height)


def fill_template_pptx(
    template_bytes: bytes,
    analyzed: AnalyzedTemplate,
    blueprint: PresentationBlueprint,
    slide_images: Dict[int, bytes],
) -> bytes:
    template_prs = Presentation(io.BytesIO(template_bytes))
    prs = Presentation(io.BytesIO(template_bytes))
    prs.slide_width = template_prs.slide_width
    prs.slide_height = template_prs.slide_height
    _clear_slides(prs)

    runtime_by_key = {rt.schema.template_key: rt for rt in analyzed.runtimes}

    for index, bp_slide in enumerate(blueprint.slides):
        runtime = runtime_by_key.get(bp_slide.template_key)
        if runtime is None:
            runtime = analyzed.runtimes[min(index, len(analyzed.runtimes) - 1)]

        src_index = runtime.schema.slide_index
        if src_index < 0 or src_index >= len(template_prs.slides):
            src_index = min(index, len(template_prs.slides) - 1)

        slide = clone_slide_from_template(template_prs, src_index, prs)
        _fill_slide(slide, bp_slide, runtime)

        if analyzed.user_style:
            apply_slide_typography(slide, user_style=analyzed.user_style)

        if index in slide_images:
            _insert_image(slide, slide_images[index], prs.slide_width, prs.slide_height)

    output = io.BytesIO()
    prs.save(output)
    built = output.getvalue()
    logger.info(
        "PPTX из шаблона: %s слайдов, %s КБ (клонирование образцов, без перезаписи theme)",
        len(blueprint.slides),
        len(built) // 1024,
    )
    return built

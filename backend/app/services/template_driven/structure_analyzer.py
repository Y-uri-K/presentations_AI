from __future__ import annotations

import io
import logging
import re
from dataclasses import dataclass
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from app.config import get_settings
from app.schemas.template_blueprint import ContentSlotSchema, TemplateCatalog, TemplateSlideSchema
from app.services.template_style_extractor import UserTemplateStyle
from app.services.pptx_kinds import classify_slide, layout_name
from app.services.template_driven.constants import MIN_WORDS_BY_ROLE

logger = logging.getLogger(__name__)
settings = get_settings()


@dataclass
class TemplateSlideRuntime:
    schema: TemplateSlideSchema
    layout: object


@dataclass
class AnalyzedTemplate:
    catalog: TemplateCatalog
    runtimes: List[TemplateSlideRuntime]
    user_style: Optional[UserTemplateStyle] = None


def _word_count(text: str) -> int:
    return len(re.findall(r"\w+", text, flags=re.UNICODE))


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text_frame.text.strip()


def _placeholder_role(ph_type, ph_idx: int) -> str:
    if ph_type == PP_PLACEHOLDER.TITLE:
        return "title"
    if ph_type == PP_PLACEHOLDER.CENTER_TITLE:
        return "title"
    if ph_type == PP_PLACEHOLDER.SUBTITLE:
        return "subtitle"
    if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
        return "body"
    return f"placeholder_{ph_idx}"


def _infer_slide_type(slide, *, card_shapes: int, has_table: bool, has_chart: bool) -> str:
    base = classify_slide(slide)
    if has_table:
        return "table"
    if has_chart:
        return "diagram"
    if card_shapes >= 4:
        return "cards"
    if card_shapes == 2:
        left = 0
        right = 0
        mid = slide.part.slide_width / 2 if hasattr(slide.part, "slide_width") else 0
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and not getattr(shape, "is_placeholder", False):
                if shape.left < mid:
                    left += 1
                else:
                    right += 1
        if left and right:
            return "comparison"
    if base == "title":
        return "title"
    if base == "section":
        return "section"
    if base == "image_content":
        return "diagram"
    if card_shapes >= 2:
        return "cards"
    return "title_content"


def _analyze_slide(slide, slide_index: int) -> TemplateSlideSchema:
    slots: List[ContentSlotSchema] = []
    card_slot_ids: List[str] = []
    column_slot_ids: List[str] = []
    metric_groups: List[List[str]] = []
    has_table = False
    has_chart = False
    has_picture = False
    card_shapes = 0

    title_slot: Optional[str] = None
    subtitle_slot: Optional[str] = None

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            has_table = True
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            has_chart = True
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            has_picture = True

    ph_idx = 0
    for shape in slide.placeholders:
        try:
            ph_type = shape.placeholder_format.type
            idx = shape.placeholder_format.idx
        except (AttributeError, ValueError):
            ph_idx += 1
            continue
        role = _placeholder_role(ph_type, idx)
        slot_id = f"ph_{idx}"
        default_text = _shape_text(shape)
        min_words = MIN_WORDS_BY_ROLE.get(role, MIN_WORDS_BY_ROLE["default"])
        slots.append(
            ContentSlotSchema(
                slot_id=slot_id,
                role=role,
                placeholder_idx=idx,
                min_words=min_words,
                default_text=default_text,
            )
        )
        if role == "title":
            title_slot = slot_id
        elif role == "subtitle":
            subtitle_slot = slot_id
        ph_idx += 1

    extra_index = 0
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        text = _shape_text(shape)
        if not text and shape.shape_type not in (MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.CHART):
            continue
        is_card_shape = shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX)
        if is_card_shape:
            card_shapes += 1
        slot_id = f"shape_{extra_index}"
        extra_index += 1
        role = "card_body" if is_card_shape else "body"
        slots.append(
            ContentSlotSchema(
                slot_id=slot_id,
                role=role,
                min_words=MIN_WORDS_BY_ROLE.get(role, 15),
                default_text=text,
            )
        )
        if role == "card_body":
            card_slot_ids.append(slot_id)

    slide_type = _infer_slide_type(
        slide, card_shapes=card_shapes, has_table=has_table, has_chart=has_chart
    )
    if slide_type == "comparison" and len(card_slot_ids) >= 2:
        column_slot_ids = card_slot_ids[:2]

    required = [s.slot_id for s in slots if s.required]

    template_key = f"{slide_type}_{slide_index}_{layout_name(slide.slide_layout)}"

    return TemplateSlideSchema(
        template_key=template_key,
        slide_index=slide_index,
        layout_name=layout_name(slide.slide_layout) or slide_type,
        slide_type=slide_type,  # type: ignore[arg-type]
        title_slot=title_slot,
        subtitle_slot=subtitle_slot,
        content_slots=slots,
        card_slots=card_slot_ids,
        column_slots=column_slot_ids,
        has_table=has_table,
        has_chart=has_chart,
        has_picture=has_picture,
        has_timeline=slide_type in ("timeline", "process"),
        has_metrics=slide_type == "kpi",
        metric_slot_groups=metric_groups,
        required_fields=required,
    )


def analyze_template_structure(template_bytes: bytes) -> AnalyzedTemplate:
    prs = Presentation(io.BytesIO(template_bytes))
    limit = min(len(prs.slides), settings.presentation_max_slides)
    schemas: List[TemplateSlideSchema] = []
    runtimes: List[TemplateSlideRuntime] = []
    seen_keys: set[str] = set()

    for index in range(limit):
        slide = prs.slides[index]
        schema = _analyze_slide(slide, index)
        if schema.template_key in seen_keys:
            schema = schema.model_copy(
                update={"template_key": f"{schema.template_key}_dup{index}"}
            )
        seen_keys.add(schema.template_key)
        schemas.append(schema)
        runtimes.append(TemplateSlideRuntime(schema=schema, layout=slide.slide_layout))

    if not schemas:
        for layout in prs.slide_layouts:
            name = layout_name(layout)
            if "blank" in name or "пуст" in name:
                continue
            schema = TemplateSlideSchema(
                template_key=f"layout_{name}",
                slide_index=-1,
                layout_name=name,
                slide_type="title_content",
                content_slots=[
                    ContentSlotSchema(slot_id="ph_0", role="title", placeholder_idx=0, min_words=5),
                    ContentSlotSchema(slot_id="ph_1", role="body", placeholder_idx=1, min_words=30),
                ],
                required_fields=["ph_0", "ph_1"],
            )
            schemas.append(schema)
            runtimes.append(TemplateSlideRuntime(schema=schema, layout=layout))
            break

    logger.info(
        "Анализ шаблона: %s слайдов-образцов, типы=%s",
        len(schemas),
        [s.slide_type for s in schemas],
    )
    return AnalyzedTemplate(catalog=TemplateCatalog(slides=schemas), runtimes=runtimes)

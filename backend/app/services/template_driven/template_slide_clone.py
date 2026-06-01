from __future__ import annotations

from copy import deepcopy

from pptx import Presentation


def clone_slide_from_template(
    source_prs: Presentation,
    source_index: int,
    dest_prs: Presentation,
):
    """Копирует слайд-образец из шаблона со всеми фигурами (не только layout)."""
    if source_index < 0 or source_index >= len(source_prs.slides):
        source_index = 0
    source = source_prs.slides[source_index]
    dest = dest_prs.slides.add_slide(source.slide_layout)
    for shape in source.shapes:
        newel = deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(newel, "p:extLst")
    return dest

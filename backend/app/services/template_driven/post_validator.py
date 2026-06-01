from __future__ import annotations

import io
import logging
from dataclasses import dataclass
from typing import List

from pptx import Presentation

from app.services.template_driven.text_utils import is_placeholder_text, word_count

logger = logging.getLogger(__name__)


@dataclass
class PptxSlideIssue:
    slide_index: int
    snippet: str
    reason: str


def validate_filled_pptx(pptx_bytes: bytes) -> List[PptxSlideIssue]:
    prs = Presentation(io.BytesIO(pptx_bytes))
    issues: List[PptxSlideIssue] = []

    for index in range(len(prs.slides)):
        slide = prs.slides[index]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if not text:
                issues.append(
                    PptxSlideIssue(slide_index=index, snippet="", reason="Пустой текстовый блок")
                )
                continue
            if is_placeholder_text(text):
                issues.append(
                    PptxSlideIssue(
                        slide_index=index,
                        snippet=text[:60],
                        reason="Placeholder-текст",
                    )
                )
            elif word_count(text) < 3 and len(text) < 15:
                issues.append(
                    PptxSlideIssue(
                        slide_index=index,
                        snippet=text[:60],
                        reason="Слишком короткий блок",
                    )
                )
    return issues

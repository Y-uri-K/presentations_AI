from __future__ import annotations

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.schemas.semantic_slides import ComparisonSlide, KpiSlide, MetricItem, ResultsSlide
from app.services.slide_renderers.context import RenderContext
from app.services.slide_renderers.drawing import add_title
from app.services.slide_renderers.layout_bounds import content_bounds_for_slide
from app.services.slide_palette_colors import palette_rgb
from app.services.slide_renderers.content_cell import (
    render_content_cell,
    render_equal_cells_row,
    render_equal_cells_stack,
)
from app.services.slide_templates._helpers import render_with_title
from app.services.slide_templates.image_layouts import (
    render_image_comparison_stack,
    render_image_metrics_column,
)
from app.services.slide_templates.visual_layouts import render_dense_kpi_grid, render_dense_kpi_row


def render_comparison(ctx: RenderContext, variant: str) -> None:
    spec: ComparisonSlide = ctx.spec  # type: ignore[assignment]
    if variant == "image_stack_top" or ctx.has_image_zone:
        render_image_comparison_stack(ctx, spec)
        return

    render_with_title(ctx)
    bounds = content_bounds_for_slide(ctx)
    span_h = bounds.bottom_pct - bounds.top_pct
    span_w = bounds.right_pct - bounds.left_pct
    gap = 0.016

    if variant == "versus_center":
        col_w = (span_w - gap) / 2
        add_comparison_columns(ctx, bounds, col_w, gap, spec)
        mid = (bounds.left_pct + bounds.right_pct) / 2
        cy = int(ctx.slide_height * (bounds.top_pct + span_h / 2))
        cx = int(ctx.slide_width * mid)
        badge = ctx.slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            cx - Emu(Inches(0.35)),
            cy - Emu(Inches(0.35)),
            Emu(Inches(0.7)),
            Emu(Inches(0.7)),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = ctx.accent
        badge.line.fill.background()
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "VS"
        r.font.bold = True
        r.font.size = Pt(14)
        r.font.color.rgb = palette_rgb(ctx.user_style, "on_accent")
    elif variant in ("stacked_rows", "mirror_bars"):
        render_equal_cells_stack(
            ctx,
            bounds,
            [
                (spec.left.heading, "\n".join(spec.left.points[:6])),
                (spec.right.heading, "\n".join(spec.right.points[:6])),
            ],
            accent_index=1,
            max_items=2,
        )
    else:
        col_w = (span_w - gap) / 2
        add_comparison_columns(ctx, bounds, col_w, gap, spec, accent_right=True)


def add_comparison_columns(
    ctx: RenderContext,
    bounds,
    col_w: float,
    gap: float,
    spec: ComparisonSlide,
    *,
    accent_right: bool = False,
) -> None:
    span_h = bounds.bottom_pct - bounds.top_pct
    render_content_cell(
        ctx,
        left_pct=bounds.left_pct,
        top_pct=bounds.top_pct,
        width_pct=col_w,
        height_pct=span_h,
        heading=spec.left.heading,
        body="\n".join(f"• {p}" for p in spec.left.points[:6]),
        accent=False,
    )
    render_content_cell(
        ctx,
        left_pct=bounds.left_pct + col_w + gap,
        top_pct=bounds.top_pct,
        width_pct=col_w,
        height_pct=span_h,
        heading=spec.right.heading,
        body="\n".join(f"• {p}" for p in spec.right.points[:6]),
        accent=accent_right,
    )


def render_kpi(ctx: RenderContext, variant: str) -> None:
    spec: KpiSlide = ctx.spec  # type: ignore[assignment]
    metrics = spec.metrics[:6]
    if variant == "image_metrics_column" or ctx.has_image_zone:
        render_image_metrics_column(ctx, metrics)
        return

    render_with_title(ctx)
    bounds = content_bounds_for_slide(ctx)
    span_h = bounds.bottom_pct - bounds.top_pct
    span_w = bounds.right_pct - bounds.left_pct

    if variant == "hero_metric" and metrics:
        first = metrics[0]
        render_content_cell(
            ctx,
            left_pct=bounds.left_pct,
            top_pct=bounds.top_pct,
            width_pct=span_w,
            height_pct=span_h * 0.34,
            heading=first.label,
            body=first.note or "",
            kpi_value=first.value,
            accent=True,
        )
        rest = metrics[1:5]
        if rest:
            items = [
                (m.label, (m.note or "").strip() or m.label, m.value) for m in rest
            ]
            render_equal_cells_row(
                ctx,
                left_pct=bounds.left_pct,
                top_pct=bounds.top_pct + span_h * 0.36 + 0.02,
                width_pct=span_w,
                height_pct=span_h * 0.64 - 0.02,
                items=items,
            )
    elif variant in ("accent_row", "kpi_trio"):
        render_dense_kpi_row(ctx, metrics[:3] if variant == "kpi_trio" else metrics)
    elif variant == "delta_cards":
        pairs = [(m.value, f"{m.label}\n{m.note or ''}".strip()) for m in metrics[:4]]
        render_equal_cells_stack(ctx, bounds, pairs, accent_index=0)
    else:
        render_dense_kpi_grid(ctx, metrics)


def render_results(ctx: RenderContext, variant: str) -> None:
    spec: ResultsSlide = ctx.spec  # type: ignore[assignment]
    if variant == "image_metrics_column" or ctx.has_image_zone:
        metrics = [
            MetricItem(value=r.value, label=r.label, note=r.trend) for r in spec.results[:6]
        ]
        render_image_metrics_column(ctx, metrics)
        return
    if spec.summary:
        add_title(ctx)
        from app.services.slide_renderers.drawing import add_subtitle

        add_subtitle(ctx, spec.summary, top_pct=0.20)
    else:
        render_with_title(ctx)
    kpi_spec = KpiSlide(
        type="kpi",
        title=spec.title,
        metrics=[
            MetricItem(value=r.value, label=r.label, note=r.trend) for r in spec.results[:6]
        ],
        image=spec.image,
        speaker_notes=spec.speaker_notes,
    )
    kpi_ctx = RenderContext(
        slide=ctx.slide,
        spec=kpi_spec,
        slide_width=ctx.slide_width,
        slide_height=ctx.slide_height,
        user_style=ctx.user_style,
        image_bytes=ctx.image_bytes,
        slide_index=ctx.slide_index,
        layout_variant=variant if variant in ("hero_metric", "accent_row") else "metric_grid",
    )
    render_kpi(kpi_ctx, kpi_ctx.layout_variant)

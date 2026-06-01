from __future__ import annotations

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu, Inches

from app.schemas.semantic_slides import ProcessSlide, TimelineSlide
from app.services.slide_renderers.context import RenderContext
from app.services.slide_renderers.drawing import (
    add_card,
    add_horizontal_timeline,
    add_process_flow,
)
from app.services.slide_renderers.layout_bounds import content_bounds_for_slide
from app.services.slide_palette_colors import palette_rgb
from app.services.slide_templates._helpers import render_with_title


def _vertical_rail(ctx: RenderContext, steps: list) -> None:
    bounds = content_bounds_for_slide(ctx)
    accent = ctx.accent
    line_left = int(ctx.slide_width * (bounds.left_pct + 0.04))
    y0 = int(ctx.slide_height * bounds.top_pct)
    y1 = int(ctx.slide_height * bounds.bottom_pct)
    line = ctx.slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        line_left,
        y0,
        Emu(Inches(0.04)),
        y1 - y0,
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()

    step_h = (bounds.bottom_pct - bounds.top_pct) / max(len(steps), 1)
    for index, (label, description) in enumerate(steps[:6]):
        top = bounds.top_pct + index * step_h
        node = ctx.slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            line_left - Emu(Inches(0.08)),
            int(ctx.slide_height * top),
            Emu(Inches(0.2)),
            Emu(Inches(0.2)),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = accent
        node.line.fill.background()
        add_card(
            ctx,
            left_pct=bounds.left_pct + 0.10,
            top_pct=top,
            width_pct=bounds.right_pct - bounds.left_pct - 0.12,
            height_pct=step_h * 0.88,
            heading=label,
            body=description or "",
            style="flat",
        )


def render_timeline(ctx: RenderContext, variant: str) -> None:
    spec: TimelineSlide = ctx.spec  # type: ignore[assignment]
    steps = [(s.label, s.description or "") for s in spec.steps]
    if variant == "image_column" or ctx.has_image_zone:
        from app.services.slide_templates.image_layouts import render_image_column_texts

        render_image_column_texts(
            ctx,
            [f"{label}: {desc}" if desc else label for label, desc in steps[:6]],
        )
        return

    render_with_title(ctx)

    if variant == "vertical_rail":
        _vertical_rail(ctx, steps)
    elif variant == "step_cards":
        bounds = content_bounds_for_slide(ctx)
        count = min(len(steps), 5)
        gap = 0.02
        w = (bounds.right_pct - bounds.left_pct - gap * (count - 1)) / max(count, 1)
        for index, (label, description) in enumerate(steps[:count]):
            add_card(
                ctx,
                left_pct=bounds.left_pct + index * (w + gap),
                top_pct=bounds.top_pct,
                width_pct=w,
                height_pct=bounds.bottom_pct - bounds.top_pct,
                heading=f"Шаг {index + 1}",
                body=f"{label}\n{description or ''}".strip(),
                accent=index == count - 1,
                style="rounded",
            )
    elif variant == "milestones":
        bounds = content_bounds_for_slide(ctx)
        count = min(len(steps), 5)
        gap = 0.02
        w = (bounds.right_pct - bounds.left_pct - gap * (count - 1)) / max(count, 1)
        for index, (label, description) in enumerate(steps[:count]):
            add_card(
                ctx,
                left_pct=bounds.left_pct + index * (w + gap),
                top_pct=bounds.top_pct,
                width_pct=w,
                height_pct=(bounds.bottom_pct - bounds.top_pct) * 0.45,
                heading=f"M{index + 1}",
                body=label,
                accent=index == count - 1,
                style="rounded",
            )
            if description:
                add_card(
                    ctx,
                    left_pct=bounds.left_pct + index * (w + gap),
                    top_pct=bounds.top_pct + (bounds.bottom_pct - bounds.top_pct) * 0.5,
                    width_pct=w,
                    height_pct=(bounds.bottom_pct - bounds.top_pct) * 0.42,
                    heading="",
                    body=description,
                    style="flat",
                )
    else:
        add_horizontal_timeline(ctx, steps)


def render_process(ctx: RenderContext, variant: str) -> None:
    spec: ProcessSlide = ctx.spec  # type: ignore[assignment]
    steps = [(s.title, s.description or "") for s in spec.steps]
    if variant == "image_column" or ctx.has_image_zone:
        from app.services.slide_templates.image_layouts import render_image_column_texts

        render_image_column_texts(
            ctx,
            [f"{title}: {desc}" if desc else title for title, desc in steps[:4]],
        )
        return

    render_with_title(ctx)

    if variant == "vertical_pipe":
        _vertical_rail(ctx, steps)
    elif variant == "step_badges":
        from app.services.slide_templates._helpers import render_sidebar_list

        render_sidebar_list(
            ctx,
            [(f"{i + 1}. {title}", description or "") for i, (title, description) in enumerate(steps[:5])],
        )
    elif variant == "funnel":
        bounds = content_bounds_for_slide(ctx)
        levels = min(len(steps), 4)
        for index, (title, description) in enumerate(steps[:levels]):
            width_frac = 0.9 - index * 0.12
            left_frac = (1.0 - width_frac) / 2
            add_card(
                ctx,
                left_pct=bounds.left_pct + left_frac * (bounds.right_pct - bounds.left_pct),
                top_pct=bounds.top_pct + index * 0.12,
                width_pct=width_frac * (bounds.right_pct - bounds.left_pct),
                height_pct=0.10,
                heading=title,
                body=description or "",
                accent=index == 0,
                style="rounded",
            )
    else:
        add_process_flow(ctx, steps)

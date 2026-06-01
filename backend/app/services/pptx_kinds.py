from __future__ import annotations

from copy import deepcopy
from dataclasses import dataclass, field
from typing import Dict, List, Optional

from pptx.enum.shapes import PP_PLACEHOLDER

BLANK_LAYOUT_KEYWORDS = ("blank", "пуст", "empty")

KIND_KEYWORDS = {
    "title": ("title slide", "титул", "титульный", "cover", "облож"),
    "section": ("section", "section header", "раздел", "divider", "заголовок раздела"),
    "image_content": (
        "picture",
        "image",
        "photo",
        "изображ",
        "фото",
        "caption",
        "media",
        "content with caption",
    ),
    "title_content": (
        "title and content",
        "content",
        "object",
        "текст",
        "body",
        "comparison",
        "two content",
        "объект",
        "заголовок и объект",
    ),
}

KIND_FALLBACK = {
    "title": ("title", "title_content", "section"),
    "section": ("section", "title_content", "title"),
    "image_content": ("image_content", "title_content"),
    "title_content": ("title_content", "image_content", "section"),
    "blank": ("blank", "title_content", "section"),
}


@dataclass
class SlideBlueprint:
    kind: str
    layout: object
    decorations: List = field(default_factory=list)
    layout_name: str = ""
    background_xml: Optional[object] = None


def layout_name(layout) -> str:
    return (layout.name or "").lower()


def is_blank_layout(layout) -> bool:
    return any(keyword in layout_name(layout) for keyword in BLANK_LAYOUT_KEYWORDS)


def classify_by_name(name: str) -> Optional[str]:
    lowered = name.lower()
    for kind, keywords in KIND_KEYWORDS.items():
        if any(keyword in lowered for keyword in keywords):
            return kind
    return None


def classify_layout(layout) -> str:
    if is_blank_layout(layout):
        return "blank"
    return classify_by_name(layout_name(layout)) or "title_content"


def classify_slide(slide) -> str:
    layout_kind = classify_layout(slide.slide_layout)
    if layout_kind != "title_content":
        return layout_kind

    placeholder_types = set()
    for shape in slide.placeholders:
        try:
            placeholder_types.add(shape.placeholder_format.type)
        except (AttributeError, ValueError):
            continue

    if PP_PLACEHOLDER.PICTURE in placeholder_types:
        return "image_content"
    if PP_PLACEHOLDER.TITLE in placeholder_types and PP_PLACEHOLDER.BODY not in placeholder_types:
        if any(k in layout_name(slide.slide_layout) for k in ("section", "раздел")):
            return "section"
        return "title"
    return "title_content"


def extract_slide_background_xml(slide):
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    c_sld = slide.element.find("p:cSld", ns)
    if c_sld is None:
        return None
    bg = c_sld.find("p:bg", ns)
    if bg is None:
        return None
    return deepcopy(bg)


def extract_blueprint(slide, kind: str) -> SlideBlueprint:
    decorations = []
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False):
            continue
        decorations.append(deepcopy(shape.element))
    return SlideBlueprint(
        kind=kind,
        layout=slide.slide_layout,
        decorations=decorations,
        layout_name=layout_name(slide.slide_layout),
        background_xml=extract_slide_background_xml(slide),
    )


def build_slide_catalog(prs) -> Dict[str, SlideBlueprint]:
    catalog: Dict[str, SlideBlueprint] = {}

    for slide in prs.slides:
        kind = classify_slide(slide)
        if kind == "blank" or kind in catalog:
            continue
        catalog[kind] = extract_blueprint(slide, kind)

    for layout in prs.slide_layouts:
        kind = classify_layout(layout)
        if kind == "blank" or kind in catalog:
            continue
        catalog[kind] = SlideBlueprint(
            kind=kind,
            layout=layout,
            decorations=[],
            layout_name=layout_name(layout),
        )

    if not catalog:
        layouts = [layout for layout in prs.slide_layouts if not is_blank_layout(layout)]
        if layouts:
            catalog["title_content"] = SlideBlueprint(
                kind="title_content",
                layout=layouts[0],
                decorations=[],
                layout_name=layout_name(layouts[0]),
            )

    return catalog


def resolve_blueprint(catalog: Dict[str, SlideBlueprint], kind: str) -> SlideBlueprint:
    for candidate in KIND_FALLBACK.get(kind, (kind,)):
        if candidate in catalog:
            return catalog[candidate]
    return next(iter(catalog.values()))


def resolve_spec_kind(spec_layout: str, has_image: bool) -> str:
    if has_image and spec_layout in ("image_content", "title_content"):
        return "image_content"
    if spec_layout in KIND_FALLBACK:
        return spec_layout
    return "title_content"


SEMANTIC_TO_LAYOUT_KIND = {
    "title": "title",
    "thank_you": "title",
    "conclusion": "section",
    "agenda": "blank",
    "problem": "blank",
    "goals": "blank",
    "kpi": "blank",
    "cards": "blank",
    "comparison": "blank",
    "timeline": "blank",
    "process": "blank",
    "results": "blank",
    "table": "blank",
    "diagram": "image_content",
}


def resolve_semantic_kind(slide_type: str, *, has_image: bool) -> str:
    kind = SEMANTIC_TO_LAYOUT_KIND.get(slide_type, "blank")
    if has_image and slide_type == "diagram":
        return "image_content"
    return kind
